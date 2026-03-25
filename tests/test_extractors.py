"""Tests for content extractors — HTML, PDF, arXiv, GitHub, API adapters, and registry routing."""

from unittest.mock import MagicMock, patch

import pytest

from mini_rag.extractors import (
    ArxivExtractor,
    DocxExtractor,
    EpubExtractor,
    GenericExtractor,
    GitHubExtractor,
    PDFExtractor,
    PptxExtractor,
    RssFeedExtractor,
    RtfExtractor,
    SpreadsheetExtractor,
    DevToFetcher,
    FandomFetcher,
    MarkdownPassthroughExtractor,
    RedditFetcher,
    ScrapedPage,
    SemanticScholarFetcher,
    StackExchangeFetcher,
    WikipediaFetcher,
    YouTubeFetcher,
    get_direct_fetcher,
    _slugify,
    extract_content,
    get_extractor,
    save_scraped_page,
)


# ─── Extractor routing ───


class TestExtractorRouting:
    def test_html_routes_to_generic(self):
        ext = get_extractor("https://example.com/page", "text/html; charset=utf-8")
        assert isinstance(ext, GenericExtractor)

    def test_pdf_routes_to_pdf(self):
        ext = get_extractor("https://example.com/doc.pdf", "application/pdf")
        assert isinstance(ext, PDFExtractor)

    def test_pdf_url_extension(self):
        ext = get_extractor("https://example.com/doc.pdf", "application/octet-stream")
        assert isinstance(ext, PDFExtractor)

    def test_arxiv_routes_to_arxiv(self):
        ext = get_extractor("https://arxiv.org/abs/2301.00001", "text/html")
        assert isinstance(ext, ArxivExtractor)

    def test_github_routes_to_github(self):
        ext = get_extractor("https://github.com/user/repo", "text/html")
        assert isinstance(ext, GitHubExtractor)

    def test_unknown_content_type(self):
        ext = get_extractor("https://example.com/file", "application/zip")
        assert ext is None

    def test_specificity_order(self):
        """arXiv PDF should route to ArxivExtractor, not PDFExtractor."""
        ext = get_extractor("https://arxiv.org/pdf/2301.00001", "application/pdf")
        assert isinstance(ext, ArxivExtractor)


# ─── GenericExtractor ───


class TestGenericExtractor:
    @pytest.fixture
    def extractor(self):
        return GenericExtractor()

    def test_can_handle_html(self, extractor):
        assert extractor.can_handle("https://example.com", "text/html")
        assert extractor.can_handle("https://example.com", "text/html; charset=utf-8")

    def test_cannot_handle_pdf(self, extractor):
        assert not extractor.can_handle("https://example.com", "application/pdf")

    def test_basic_extraction(self, extractor):
        html = b"""<html><head><title>Test Page</title></head>
        <body>
        <article>
        <h1>Main Heading</h1>
        <p>This is a paragraph with enough content to pass the minimum threshold.
        It contains multiple sentences to ensure adequate word count for extraction.
        The content should be properly converted to markdown format.</p>
        <h2>Sub Heading</h2>
        <p>Another paragraph with additional content that adds to the word count.</p>
        </article>
        </body></html>"""
        page = extractor.extract("https://example.com/test", html, "text/html")
        assert page is not None
        assert page.title == "Test Page"
        assert page.word_count > 0
        assert "Main Heading" in page.content
        assert page.source_type == "web"

    def test_strips_scripts_and_styles(self, extractor):
        html = b"""<html><head><title>Test</title></head>
        <body>
        <article>
        <p>Real content that should be extracted and preserved in output.
        This paragraph needs enough text to pass the minimum content check.</p>
        <script>alert('evil');</script>
        <style>.hidden { display: none; }</style>
        <p>More real content after the script and style tags are stripped.</p>
        </article>
        </body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert "alert" not in page.content
        assert "display: none" not in page.content

    def test_strips_nav_footer(self, extractor):
        html = b"""<html><head><title>Test</title></head>
        <body>
        <nav><a href="/">Home</a><a href="/about">About</a></nav>
        <article>
        <p>Main content that matters. This is what should be extracted from the page.
        Multiple sentences here to ensure we pass the minimum threshold.</p>
        </article>
        <footer>Copyright 2024</footer>
        </body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert "Copyright" not in page.content

    def test_extracts_links(self, extractor):
        html = b"""<html><head><title>Links</title></head>
        <body>
        <article>
        <p>Content with links and enough text to pass minimum threshold easily.
        This paragraph provides sufficient word count for extraction.</p>
        <a href="https://other.com/page">External Link</a>
        <a href="#section">Anchor</a>
        <a href="javascript:void(0)">JS</a>
        </article>
        </body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert "https://other.com/page" in page.links
        # Anchors and javascript links should be excluded
        assert not any("#section" in l for l in page.links)

    def test_title_from_og(self, extractor):
        html = b"""<html><head>
        <title>Site Name - Long Title</title>
        <meta property="og:title" content="Clean Title" />
        </head><body>
        <article><p>Content body with enough words to pass the minimum extraction threshold.
        This ensures the page is not filtered out as too thin.</p></article>
        </body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert page.title == "Clean Title"

    def test_title_strips_suffix(self, extractor):
        html = b"""<html><head>
        <title>Article Title | Site Name</title>
        </head><body>
        <article><p>Content body with sufficient text to pass extraction minimum.
        Multiple sentences ensure adequate word count.</p></article>
        </body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert page.title == "Article Title"

    def test_insufficient_content_returns_none(self, extractor):
        html = b"<html><body><p>Hi</p></body></html>"
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is None

    def test_table_to_markdown(self, extractor):
        html = b"""<html><head><title>Table</title></head>
        <body><article>
        <p>Intro text with enough content to pass minimum threshold requirements.
        The table below shows structured data in markdown format.</p>
        <table>
        <tr><th>Name</th><th>Value</th></tr>
        <tr><td>Alpha</td><td>100</td></tr>
        <tr><td>Beta</td><td>200</td></tr>
        </table>
        </article></body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert "| Name | Value |" in page.content
        assert "| Alpha | 100 |" in page.content

    def test_code_block_extraction(self, extractor):
        html = b"""<html><head><title>Code</title></head>
        <body><article>
        <p>Some explanatory text about the code example below, with enough words
        to pass the minimum content threshold for extraction.</p>
        <pre><code class="language-python">def hello():
    print("world")</code></pre>
        </article></body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert "```python" in page.content
        assert 'print("world")' in page.content

    def test_content_selector_priority(self, extractor):
        html = b"""<html><head><title>Test</title></head>
        <body>
        <div class="sidebar">Sidebar noise that should not be the main content.</div>
        <article>
        <p>Main article content that should be identified as the primary content area.
        This has substantially more text than the sidebar.</p>
        </article>
        <div class="ads">Advertisement content that is also noise.</div>
        </body></html>"""
        page = extractor.extract("https://example.com", html, "text/html")
        assert page is not None
        assert "Main article" in page.content


# ─── GitHubExtractor ───


class TestGitHubExtractor:
    @pytest.fixture
    def extractor(self):
        return GitHubExtractor()

    def test_can_handle(self, extractor):
        assert extractor.can_handle("https://github.com/user/repo", "text/html")
        assert not extractor.can_handle("https://github.com/user/repo", "application/json")
        assert not extractor.can_handle("https://example.com", "text/html")

    def test_repo_page_with_readme(self, extractor):
        html = b"""<html><head><title>user/repo</title></head>
        <body>
        <p class="f4 my-3">A cool project description</p>
        <article class="markdown-body">
        <h1>My Project</h1>
        <p>This is the README content with installation instructions and usage details.
        It provides comprehensive documentation for the repository.</p>
        </article>
        </body></html>"""
        page = extractor.extract("https://github.com/user/repo", html, "text/html")
        assert page is not None
        assert page.source_type == "github"
        assert "user/repo" in page.title

    def test_non_repo_page_returns_none(self, extractor):
        """Non-repo GitHub pages (trending, explore) should return None."""
        html = b"""<html><head><title>Trending</title></head>
        <body><div>Some trending content</div></body></html>"""
        page = extractor.extract("https://github.com/trending", html, "text/html")
        assert page is None


# ─── ArxivExtractor ───


class TestArxivExtractor:
    @pytest.fixture
    def extractor(self):
        return ArxivExtractor()

    def test_can_handle(self, extractor):
        assert extractor.can_handle("https://arxiv.org/abs/2301.00001", "text/html")
        assert extractor.can_handle("https://arxiv.org/pdf/2301.00001", "application/pdf")
        assert not extractor.can_handle("https://example.com", "text/html")

    def test_abstract_page(self, extractor):
        html = b"""<html><body>
        <h1 class="title">Title:Quantum Computing Survey</h1>
        <div class="authors">Authors:Alice, Bob</div>
        <blockquote class="abstract">Abstract:This paper surveys quantum computing advances.</blockquote>
        <div class="subjects">Subjects: Quantum Physics</div>
        <div class="dateline">Submitted 2024-01-15</div>
        <a href="/pdf/2301.00001">PDF</a>
        </body></html>"""
        page = extractor.extract("https://arxiv.org/abs/2301.00001", html, "text/html")
        assert page is not None
        assert page.source_type == "arxiv"
        assert "Quantum Computing Survey" in page.title
        assert "Alice, Bob" in page.content
        assert "quantum computing" in page.content.lower()


# ─── ScrapedPage and save ───


class TestScrapedPage:
    def test_dataclass_defaults(self):
        page = ScrapedPage(
            url="https://example.com",
            title="Test",
            content="Hello world",
            scraped_at="2024-01-01T00:00:00",
            word_count=2,
        )
        assert page.links == []
        assert page.source_type == "web"
        assert page.raw_bytes is None


class TestSaveScrapedPage:
    def test_save_creates_file(self, tmp_path):
        page = ScrapedPage(
            url="https://example.com",
            title="Test Page",
            content="Content here",
            scraped_at="2024-01-01T00:00:00",
            word_count=2,
        )
        filepath = save_scraped_page(page, tmp_path)
        assert filepath.exists()
        assert filepath.suffix == ".md"
        text = filepath.read_text()
        assert "---" in text  # frontmatter
        assert "Test Page" in text
        assert "source_url:" in text

    def test_collision_handling(self, tmp_path):
        page = ScrapedPage(
            url="https://example.com",
            title="Same Title",
            content="Content",
            scraped_at="2024-01-01T00:00:00",
            word_count=1,
        )
        path1 = save_scraped_page(page, tmp_path)
        path2 = save_scraped_page(page, tmp_path)
        assert path1 != path2
        assert path1.exists()
        assert path2.exists()

    def test_slugify(self):
        assert _slugify("Hello World!") == "hello-world"
        assert _slugify("  Spaces  & Symbols!! ") == "spaces-symbols"
        assert len(_slugify("x" * 100)) <= 60


# ─── extract_content integration ───


class TestExtractContent:
    def test_routes_and_extracts_html(self, monkeypatch):
        """Redirect registry writes to avoid touching real scrape log."""
        import tempfile
        from pathlib import Path
        td = Path(tempfile.mkdtemp())
        monkeypatch.setattr("mini_rag.scrape_registry.SCRAPE_LOG_DIR", td)
        monkeypatch.setattr("mini_rag.scrape_registry.SCRAPE_LOG_FILE", td / "scrape-log.jsonl")

        html = b"""<html><head><title>Integration Test</title></head>
        <body><article>
        <p>Enough content here to pass all minimum thresholds for extraction.
        Multiple sentences with real words to reach the word count limit.</p>
        </article></body></html>"""
        page = extract_content("https://example.com/test", html, "text/html")
        assert page is not None
        assert page.title == "Integration Test"

        # Check scrape log was written
        log_file = td / "scrape-log.jsonl"
        assert log_file.exists()

    def test_unknown_type_logs_failure(self, monkeypatch):
        import tempfile
        from pathlib import Path
        td = Path(tempfile.mkdtemp())
        monkeypatch.setattr("mini_rag.scrape_registry.SCRAPE_LOG_DIR", td)
        monkeypatch.setattr("mini_rag.scrape_registry.SCRAPE_LOG_FILE", td / "scrape-log.jsonl")

        result = extract_content("https://example.com/file.bin", b"\x00\x01", "application/octet-stream")
        assert result is None

        import json
        log_file = td / "scrape-log.jsonl"
        if log_file.exists():
            entry = json.loads(log_file.read_text().strip())
            assert entry["success"] is False
            assert entry["extractor"] == "none"


# ─── Direct fetcher routing ───


class TestDirectFetcherRouting:
    def test_wikipedia(self):
        f = get_direct_fetcher("https://en.wikipedia.org/wiki/Python")
        assert isinstance(f, WikipediaFetcher)

    def test_wikipedia_other_lang(self):
        f = get_direct_fetcher("https://de.wikipedia.org/wiki/Python")
        assert isinstance(f, WikipediaFetcher)

    def test_youtube(self):
        f = get_direct_fetcher("https://www.youtube.com/watch?v=abc123")
        assert isinstance(f, YouTubeFetcher)

    def test_youtube_short(self):
        f = get_direct_fetcher("https://youtu.be/abc123")
        assert isinstance(f, YouTubeFetcher)

    def test_stackoverflow(self):
        f = get_direct_fetcher("https://stackoverflow.com/questions/12345/test")
        assert isinstance(f, StackExchangeFetcher)

    def test_stackexchange_subdomain(self):
        f = get_direct_fetcher("https://gaming.stackexchange.com/questions/1/test")
        assert isinstance(f, StackExchangeFetcher)

    def test_askubuntu(self):
        f = get_direct_fetcher("https://askubuntu.com/questions/99/test")
        assert isinstance(f, StackExchangeFetcher)

    def test_superuser(self):
        f = get_direct_fetcher("https://superuser.com/questions/55/test")
        assert isinstance(f, StackExchangeFetcher)

    def test_reddit(self):
        f = get_direct_fetcher("https://www.reddit.com/r/python/comments/abc/test")
        assert isinstance(f, RedditFetcher)

    def test_old_reddit(self):
        f = get_direct_fetcher("https://old.reddit.com/r/python/comments/abc/test")
        assert isinstance(f, RedditFetcher)

    def test_normal_site_returns_none(self):
        f = get_direct_fetcher("https://example.com/page")
        assert f is None


# ─── WikipediaFetcher ───


class TestWikipediaFetcher:
    @pytest.fixture
    def fetcher(self):
        return WikipediaFetcher()

    def test_can_handle(self, fetcher):
        assert fetcher.can_handle_url("https://en.wikipedia.org/wiki/Test")
        assert fetcher.can_handle_url("https://de.wikipedia.org/wiki/Test")
        assert not fetcher.can_handle_url("https://example.com/wiki/Test")

    def test_rejects_non_article_urls(self, fetcher):
        # /wiki/ path required
        page = fetcher.fetch_and_extract("https://en.wikipedia.org/w/index.php")
        assert page is None

    def test_successful_fetch(self, fetcher):
        mock_summary = MagicMock()
        mock_summary.json.return_value = {
            "title": "Python (programming language)",
            "extract": "Python is a high-level programming language.",
            "description": "General-purpose programming language",
        }
        mock_summary.raise_for_status = MagicMock()

        mock_full = MagicMock()
        mock_full.status_code = 200
        mock_full.raise_for_status = MagicMock()
        mock_full.content = b"<html><body><section><p>Full article content here.</p></section></body></html>"

        def mock_get(url, **kwargs):
            if "/summary/" in url:
                return mock_summary
            return mock_full

        with patch("mini_rag.extractors.requests.get", side_effect=mock_get):
            page = fetcher.fetch_and_extract("https://en.wikipedia.org/wiki/Python_(programming_language)")

        assert page is not None
        assert page.title == "Python (programming language)"
        assert page.source_type == "wikipedia"
        assert page.word_count > 0


# ─── YouTubeFetcher ───


class TestYouTubeFetcher:
    @pytest.fixture
    def fetcher(self):
        return YouTubeFetcher()

    def test_can_handle(self, fetcher):
        assert fetcher.can_handle_url("https://www.youtube.com/watch?v=abc")
        assert fetcher.can_handle_url("https://youtu.be/abc")
        assert not fetcher.can_handle_url("https://vimeo.com/123")

    def test_successful_fetch(self, fetcher):
        mock_oembed = MagicMock()
        mock_oembed.json.return_value = {
            "title": "Test Video Title",
            "author_name": "TestChannel",
            "author_url": "https://www.youtube.com/@TestChannel",
        }
        mock_oembed.raise_for_status = MagicMock()

        mock_page = MagicMock()
        mock_page.raise_for_status = MagicMock()
        mock_page.text = '<meta name="description" content="Video description here">'

        def mock_get(url, **kwargs):
            if "oembed" in url:
                return mock_oembed
            return mock_page

        with patch("mini_rag.extractors.requests.get", side_effect=mock_get):
            page = fetcher.fetch_and_extract("https://www.youtube.com/watch?v=test123")

        assert page is not None
        assert page.title == "Test Video Title"
        assert page.source_type == "youtube"
        assert "TestChannel" in page.content
        assert "Video description here" in page.content

    def test_short_url_normalized(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.json.return_value = {"title": "Short URL Video", "author_name": "Ch"}
        mock_resp.raise_for_status = MagicMock()
        mock_resp.text = ""

        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://youtu.be/abc123")

        assert page is not None
        assert "youtube.com/watch?v=abc123" in page.url


# ─── StackExchangeFetcher ───


class TestStackExchangeFetcher:
    @pytest.fixture
    def fetcher(self):
        return StackExchangeFetcher()

    def test_can_handle(self, fetcher):
        assert fetcher.can_handle_url("https://stackoverflow.com/questions/123/test")
        assert fetcher.can_handle_url("https://askubuntu.com/questions/1/test")
        assert fetcher.can_handle_url("https://gaming.stackexchange.com/questions/1/t")
        assert not fetcher.can_handle_url("https://example.com/questions/1/t")

    def test_rejects_non_question_urls(self, fetcher):
        page = fetcher.fetch_and_extract("https://stackoverflow.com/tags")
        assert page is None

    def test_site_param_mapping(self, fetcher):
        assert fetcher._get_site_param("stackoverflow.com") == "stackoverflow"
        assert fetcher._get_site_param("www.stackoverflow.com") == "stackoverflow"
        assert fetcher._get_site_param("gaming.stackexchange.com") == "gaming"
        assert fetcher._get_site_param("askubuntu.com") == "askubuntu"

    def test_successful_fetch(self, fetcher):
        mock_q_resp = MagicMock()
        mock_q_resp.json.return_value = {
            "items": [{
                "title": "How to test?",
                "body": "<p>Question body here</p>",
                "tags": ["python", "testing"],
                "score": 42,
                "answer_count": 3,
            }]
        }
        mock_q_resp.raise_for_status = MagicMock()

        mock_a_resp = MagicMock()
        mock_a_resp.json.return_value = {
            "items": [{
                "body": "<p>Answer body here</p>",
                "score": 15,
                "is_accepted": True,
            }]
        }
        mock_a_resp.raise_for_status = MagicMock()

        call_count = 0

        def mock_get(url, **kwargs):
            nonlocal call_count
            call_count += 1
            if "/answers" in url:
                return mock_a_resp
            return mock_q_resp

        with patch("mini_rag.extractors.requests.get", side_effect=mock_get):
            page = fetcher.fetch_and_extract("https://stackoverflow.com/questions/12345/how-to-test")

        assert page is not None
        assert page.title == "How to test?"
        assert page.source_type == "stackexchange"
        assert "python" in page.content
        assert "Answer 1" in page.content
        assert "(Accepted)" in page.content


# ─── RedditFetcher ───


class TestRedditFetcher:
    @pytest.fixture
    def fetcher(self):
        return RedditFetcher()

    def test_can_handle(self, fetcher):
        assert fetcher.can_handle_url("https://www.reddit.com/r/python/comments/abc/test")
        assert fetcher.can_handle_url("https://old.reddit.com/r/python/comments/abc/test")
        assert not fetcher.can_handle_url("https://example.com/r/python")

    def test_successful_fetch(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.json.return_value = [
            {
                "data": {
                    "children": [{
                        "data": {
                            "title": "Check out this Python trick",
                            "selftext": "Here is a cool trick using list comprehensions.",
                            "subreddit_name_prefixed": "r/python",
                            "score": 250,
                            "author": "pythonista",
                            "num_comments": 45,
                            "url": "https://www.reddit.com/r/python/comments/abc/test",
                        }
                    }]
                }
            },
            {
                "data": {
                    "children": [
                        {
                            "kind": "t1",
                            "data": {
                                "author": "commenter1",
                                "body": "Great tip! I use this all the time.",
                                "score": 30,
                            }
                        }
                    ]
                }
            },
        ]
        mock_resp.raise_for_status = MagicMock()

        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://www.reddit.com/r/python/comments/abc/test")

        assert page is not None
        assert page.title == "Check out this Python trick"
        assert page.source_type == "reddit"
        assert "r/python" in page.content
        assert "list comprehensions" in page.content
        assert "commenter1" in page.content

    def test_api_error_returns_none(self, fetcher):
        with patch("mini_rag.extractors.requests.get", side_effect=Exception("Network error")):
            page = fetcher.fetch_and_extract("https://www.reddit.com/r/test/comments/abc/test")
        assert page is None


# ─── ArxivExtractor fixes ───


class TestArxivExtractorFixes:
    @pytest.fixture
    def extractor(self):
        return ArxivExtractor()

    def test_rejects_listing_pages(self, extractor):
        """Listing pages like /list/cs.MA/ should not be handled."""
        assert not extractor.can_handle("https://arxiv.org/list/cs.MA/2026", "text/html")

    def test_handles_abs_pages(self, extractor):
        assert extractor.can_handle("https://arxiv.org/abs/2301.00001", "text/html")

    def test_handles_html_pages(self, extractor):
        assert extractor.can_handle("https://arxiv.org/html/2603.03116v1", "text/html")

    def test_handles_pdf_pages(self, extractor):
        assert extractor.can_handle("https://arxiv.org/pdf/2301.00001", "application/pdf")

    def test_html_article_extraction(self, extractor):
        html = b"""<html><body>
        <h1 class="ltx_title">A Novel Approach to Testing</h1>
        <div class="ltx_authors">Alice Smith, Bob Jones</div>
        <div class="ltx_abstract">Abstract. This paper presents a novel testing framework.</div>
        <div class="ltx_page_content">
        <section><h2>Introduction</h2><p>Testing is important for software quality.
        This section describes our approach in detail with examples.</p></section>
        <section><h2>Methods</h2><p>We used a combination of unit and integration tests
        to validate our framework across multiple domains.</p></section>
        </div>
        </body></html>"""
        page = extractor.extract("https://arxiv.org/html/2603.03116v1", html, "text/html")
        assert page is not None
        assert page.source_type == "arxiv"
        assert "Alice Smith" in page.content
        assert "novel testing" in page.content


# ─── GitHub extractor fixes ───


class TestGitHubExtractorFixes:
    @pytest.fixture
    def extractor(self):
        return GitHubExtractor()

    def test_rejects_trending(self, extractor):
        assert not extractor.can_handle("https://github.com/trending", "text/html")

    def test_rejects_explore(self, extractor):
        assert not extractor.can_handle("https://github.com/explore", "text/html")

    def test_rejects_root(self, extractor):
        assert not extractor.can_handle("https://github.com/", "text/html")

    def test_accepts_repo(self, extractor):
        assert extractor.can_handle("https://github.com/user/repo", "text/html")

    def test_accepts_repo_subpath(self, extractor):
        assert extractor.can_handle("https://github.com/user/repo/blob/main/file.py", "text/html")


# ─── MarkdownPassthroughExtractor ───


class TestMarkdownPassthroughExtractor:
    @pytest.fixture
    def extractor(self):
        return MarkdownPassthroughExtractor()

    def test_can_handle_markdown_content_type(self, extractor):
        assert extractor.can_handle("https://x.com/f", "text/markdown")
        assert extractor.can_handle("https://x.com/f", "text/markdown; charset=utf-8")
        assert extractor.can_handle("https://x.com/f", "text/x-markdown")

    def test_can_handle_plain_text(self, extractor):
        assert extractor.can_handle("https://x.com/f", "text/plain")

    def test_can_handle_md_extension(self, extractor):
        assert extractor.can_handle("https://x.com/README.md", "application/octet-stream")

    def test_can_handle_txt_extension(self, extractor):
        assert extractor.can_handle("https://x.com/notes.txt", "application/octet-stream")

    def test_does_not_handle_html(self, extractor):
        assert not extractor.can_handle("https://x.com/page", "text/html")

    def test_extract_with_heading_title(self, extractor):
        md = b"# My Document\n\nThis is the body content with enough words to pass.\n"
        page = extractor.extract("https://x.com/doc.md", md, "text/markdown")
        assert page is not None
        assert page.title == "My Document"
        assert page.source_type == "markdown"

    def test_extract_title_from_first_line(self, extractor):
        md = b"Some Document Title\n\nBody content with sufficient words here.\n"
        page = extractor.extract("https://x.com/doc.md", md, "text/markdown")
        assert page is not None
        assert page.title == "Some Document Title"

    def test_extract_title_from_filename(self, extractor):
        # Only whitespace lines before content — too long for title, falls back to filename
        long_line = "x" * 250
        md = f"\n\n\n{long_line}\n".encode()
        page = extractor.extract("https://x.com/my-notes.md", md, "text/markdown")
        assert page is not None
        assert page.title == "my-notes"

    def test_too_short_returns_none(self, extractor):
        page = extractor.extract("https://x.com/f.md", b"Hi", "text/markdown")
        assert page is None

    def test_content_passed_through(self, extractor):
        md = b"# Title\n\n```python\ndef hello():\n    pass\n```\n\nSome text here.\n"
        page = extractor.extract("https://x.com/f.md", md, "text/markdown")
        assert page is not None
        assert "```python" in page.content
        assert "def hello():" in page.content


# ─── DevToFetcher ───


class TestDevToFetcher:
    @pytest.fixture
    def fetcher(self):
        return DevToFetcher()

    def test_can_handle(self, fetcher):
        assert fetcher.can_handle_url("https://dev.to/user/article-slug")
        assert fetcher.can_handle_url("https://www.dev.to/user/article")
        assert not fetcher.can_handle_url("https://example.com/user/article")

    def test_rejects_profile_urls(self, fetcher):
        page = fetcher.fetch_and_extract("https://dev.to/username")
        assert page is None

    def test_rejects_tag_urls(self, fetcher):
        page = fetcher.fetch_and_extract("https://dev.to/t/python")
        assert page is None

    def test_successful_fetch(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.json.return_value = {
            "title": "Building a RAG System",
            "body_markdown": "# RAG Systems\n\nHere is how to build one.",
            "user": {"name": "DevAuthor"},
            "tag_list": ["python", "ai"],
            "positive_reactions_count": 42,
            "readable_publish_date": "Mar 25",
        }
        mock_resp.raise_for_status = MagicMock()

        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://dev.to/devauthor/building-a-rag-system")

        assert page is not None
        assert page.title == "Building a RAG System"
        assert page.source_type == "devto"
        assert "DevAuthor" in page.content
        assert "python" in page.content
        assert "RAG Systems" in page.content

    def test_api_error_returns_none(self, fetcher):
        with patch("mini_rag.extractors.requests.get", side_effect=Exception("API down")):
            page = fetcher.fetch_and_extract("https://dev.to/user/article")
        assert page is None

    def test_empty_response_returns_none(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.json.return_value = {"title": "", "body_markdown": ""}
        mock_resp.raise_for_status = MagicMock()
        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://dev.to/user/nonexistent")
        assert page is None


# ─── FandomFetcher ───


class TestFandomFetcher:
    @pytest.fixture
    def fetcher(self):
        return FandomFetcher()

    def test_can_handle(self, fetcher):
        assert fetcher.can_handle_url("https://callofduty.fandom.com/wiki/M4")
        assert fetcher.can_handle_url("https://starwars.fandom.com/wiki/Darth_Vader")
        assert not fetcher.can_handle_url("https://fandom.com/")
        assert not fetcher.can_handle_url("https://example.com/wiki/Test")

    def test_rejects_non_wiki_urls(self, fetcher):
        page = fetcher.fetch_and_extract("https://callofduty.fandom.com/f/community")
        assert page is None

    def test_successful_fetch(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.json.return_value = {
            "parse": {
                "title": "M4A1",
                "displaytitle": "M4A1",
                "wikitext": {"*": (
                    "The '''M4A1''' is an [[assault rifle]].\n\n"
                    "== Overview ==\nThe M4A1 is a versatile weapon.\n\n"
                    "== Variants ==\n* Standard\n* Tactical\n"
                )},
                "categories": [
                    {"*": "Assault_Rifles", "hidden": ""},
                    {"*": "Modern_Warfare"},
                ],
            }
        }
        mock_resp.raise_for_status = MagicMock()

        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://callofduty.fandom.com/wiki/M4A1")

        assert page is not None
        assert page.title == "M4A1"
        assert page.source_type == "fandom"
        assert "**M4A1**" in page.content
        assert "## Overview" in page.content
        assert "Modern Warfare" in page.content

    def test_wikitext_conversion(self):
        wikitext = (
            "'''Bold''' and ''italic'' text.\n"
            "== Section ==\n"
            "=== Subsection ===\n"
            "[[Link|Display Text]]\n"
            "[[Simple Link]]\n"
            "{{Infobox|param=value}}\n"
            "<ref>Citation</ref>\n"
        )
        md = FandomFetcher._wikitext_to_markdown(wikitext)
        assert "**Bold**" in md
        assert "*italic*" in md
        assert "## Section" in md
        assert "### Subsection" in md
        assert "Display Text" in md
        assert "Simple Link" in md
        assert "Infobox" not in md
        assert "Citation" not in md

    def test_api_error_returns_none(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.json.return_value = {"error": {"code": "missingtitle"}}
        mock_resp.raise_for_status = MagicMock()
        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://callofduty.fandom.com/wiki/Nonexistent")
        assert page is None


# ─── SemanticScholarFetcher ───


class TestSemanticScholarFetcher:
    @pytest.fixture
    def fetcher(self):
        return SemanticScholarFetcher()

    # Routing tests
    def test_can_handle_doi(self, fetcher):
        assert fetcher.can_handle_url("https://doi.org/10.1038/nature12373")
        assert fetcher.can_handle_url("https://dx.doi.org/10.1234/test")

    def test_can_handle_s2(self, fetcher):
        assert fetcher.can_handle_url("https://www.semanticscholar.org/paper/Title/abc123")

    def test_can_handle_pubmed(self, fetcher):
        assert fetcher.can_handle_url("https://pubmed.ncbi.nlm.nih.gov/12345678")

    def test_can_handle_academic_domains(self, fetcher):
        assert fetcher.can_handle_url("https://dl.acm.org/doi/10.1145/123")
        assert fetcher.can_handle_url("https://ieeexplore.ieee.org/document/123")
        assert fetcher.can_handle_url("https://www.sciencedirect.com/science/article/pii/123")
        assert fetcher.can_handle_url("https://www.researchgate.net/publication/123")
        assert fetcher.can_handle_url("https://link.springer.com/article/10.1007/123")
        assert fetcher.can_handle_url("https://www.nature.com/articles/s41586-023-123")
        assert fetcher.can_handle_url("https://journals.plos.org/plosone/article?id=123")

    def test_does_not_handle_non_academic(self, fetcher):
        assert not fetcher.can_handle_url("https://example.com/paper")
        assert not fetcher.can_handle_url("https://arxiv.org/abs/2301.00001")  # ArxivExtractor handles this

    # Paper ID resolution
    def test_resolve_doi(self, fetcher):
        assert fetcher._resolve_paper_id("https://doi.org/10.1038/nature12373") == "DOI:10.1038/nature12373"

    def test_resolve_pmid(self, fetcher):
        assert fetcher._resolve_paper_id("https://pubmed.ncbi.nlm.nih.gov/23903654") == "PMID:23903654"

    def test_resolve_s2_id(self, fetcher):
        assert fetcher._resolve_paper_id("https://semanticscholar.org/paper/Title/abc123def") == "abc123def"

    def test_resolve_url_fallback(self, fetcher):
        result = fetcher._resolve_paper_id("https://ieeexplore.ieee.org/document/12345")
        assert result.startswith("URL:")

    def test_resolve_s2_non_paper_returns_none(self, fetcher):
        assert fetcher._resolve_paper_id("https://semanticscholar.org/search") is None

    # Happy path
    def test_successful_fetch(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.status_code = 200
        mock_resp.json.return_value = {
            "title": "Attention Is All You Need",
            "abstract": "The dominant sequence transduction models are based on complex recurrent or convolutional neural networks.",
            "authors": [
                {"name": "Ashish Vaswani"},
                {"name": "Noam Shazeer"},
            ],
            "year": 2017,
            "venue": "NeurIPS",
            "citationCount": 100000,
            "referenceCount": 42,
            "tldr": {"text": "A new network architecture based solely on attention mechanisms."},
            "externalIds": {"DOI": "10.48550/arXiv.1706.03762", "ArXiv": "1706.03762"},
            "openAccessPdf": {"url": "https://arxiv.org/pdf/1706.03762"},
        }
        mock_resp.raise_for_status = MagicMock()

        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://doi.org/10.48550/arXiv.1706.03762")

        assert page is not None
        assert page.title == "Attention Is All You Need"
        assert page.source_type == "academic"
        assert "Ashish Vaswani" in page.content
        assert "NeurIPS" in page.content
        assert "100000" in page.content
        assert "TL;DR" in page.content
        assert "attention mechanisms" in page.content
        assert "transduction models" in page.content
        assert "arxiv.org/pdf" in page.content

    # Error handling
    def test_404_returns_none(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.status_code = 404
        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://doi.org/10.9999/nonexistent")
        assert page is None

    def test_api_error_returns_none(self, fetcher):
        with patch("mini_rag.extractors.requests.get", side_effect=Exception("Network error")):
            page = fetcher.fetch_and_extract("https://doi.org/10.1234/test")
        assert page is None

    def test_no_title_returns_none(self, fetcher):
        mock_resp = MagicMock()
        mock_resp.status_code = 200
        mock_resp.json.return_value = {"title": "", "abstract": "Some text"}
        mock_resp.raise_for_status = MagicMock()
        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://doi.org/10.1234/test")
        assert page is None

    def test_minimal_paper_metadata(self, fetcher):
        """Paper with only title and abstract, no TLDR or PDF."""
        mock_resp = MagicMock()
        mock_resp.status_code = 200
        mock_resp.json.return_value = {
            "title": "Old Paper",
            "abstract": "This is a historic paper from the archive.",
            "authors": [],
            "year": 1995,
            "venue": "",
            "citationCount": 0,
            "referenceCount": 0,
            "tldr": None,
            "externalIds": {},
            "openAccessPdf": None,
        }
        mock_resp.raise_for_status = MagicMock()

        with patch("mini_rag.extractors.requests.get", return_value=mock_resp):
            page = fetcher.fetch_and_extract("https://doi.org/10.1234/old")

        assert page is not None
        assert page.title == "Old Paper"
        assert "historic paper" in page.content


# ─── Direct fetcher routing for new adapters ───


class TestNewDirectFetcherRouting:
    def test_devto_routing(self):
        f = get_direct_fetcher("https://dev.to/user/article")
        assert isinstance(f, DevToFetcher)

    def test_fandom_routing(self):
        f = get_direct_fetcher("https://starwars.fandom.com/wiki/Yoda")
        assert isinstance(f, FandomFetcher)

    def test_doi_routing(self):
        f = get_direct_fetcher("https://doi.org/10.1234/test")
        assert isinstance(f, SemanticScholarFetcher)

    def test_researchgate_routing(self):
        f = get_direct_fetcher("https://www.researchgate.net/publication/123")
        assert isinstance(f, SemanticScholarFetcher)

    def test_markdown_extractor_in_registry(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/f.md", "text/markdown; charset=utf-8")
        assert isinstance(ext, MarkdownPassthroughExtractor)


# ─── Document extractor routing ───


class TestDocumentExtractorRouting:
    def test_docx_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/doc.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert isinstance(ext, DocxExtractor)

    def test_docx_by_extension(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/doc.docx", "application/octet-stream")
        assert isinstance(ext, DocxExtractor)

    def test_xlsx_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        assert isinstance(ext, SpreadsheetExtractor)

    def test_csv_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/data.csv", "text/csv")
        assert isinstance(ext, SpreadsheetExtractor)

    def test_pptx_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/slides.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        assert isinstance(ext, PptxExtractor)

    def test_epub_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/book.epub", "application/epub+zip")
        assert isinstance(ext, EpubExtractor)

    def test_rtf_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/doc.rtf", "application/rtf")
        assert isinstance(ext, RtfExtractor)

    def test_rss_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/feed.rss", "application/rss+xml")
        assert isinstance(ext, RssFeedExtractor)

    def test_atom_routing(self):
        from mini_rag.extractors import get_extractor
        ext = get_extractor("https://x.com/feed", "application/atom+xml")
        assert isinstance(ext, RssFeedExtractor)


# ─── DocxExtractor ───


class TestDocxExtractor:
    @pytest.fixture
    def extractor(self):
        return DocxExtractor()

    def test_can_handle(self, extractor):
        assert extractor.can_handle("x.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert extractor.can_handle("x.docx", "application/octet-stream")
        assert not extractor.can_handle("x.pdf", "application/pdf")

    def test_extract_real_docx(self, extractor):
        """Create a real DOCX in memory and extract it."""
        from docx import Document
        import io

        doc = Document()
        doc.core_properties.author = "Test Author"
        doc.core_properties.title = "Test Document"
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("This is the first paragraph with some test content.")
        doc.add_heading("Details", level=2)
        doc.add_paragraph("More detailed information about the topic goes here.")

        buf = io.BytesIO()
        doc.save(buf)
        raw = buf.getvalue()

        page = extractor.extract("https://example.com/test.docx", raw, "application/docx")
        assert page is not None
        assert page.source_type == "docx"
        assert page.metadata.get("author") == "Test Author"
        assert "Introduction" in page.content
        assert "first paragraph" in page.content
        assert page.word_count > 0


# ─── SpreadsheetExtractor ───


class TestSpreadsheetExtractor:
    @pytest.fixture
    def extractor(self):
        return SpreadsheetExtractor()

    def test_can_handle_csv(self, extractor):
        assert extractor.can_handle("data.csv", "text/csv")
        assert extractor.can_handle("data.tsv", "text/plain")

    def test_can_handle_xlsx(self, extractor):
        assert extractor.can_handle("data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    def test_extract_csv(self, extractor):
        csv_data = b"Name,Age,City\nAlice,30,Sydney\nBob,25,Melbourne\n"
        page = extractor.extract("https://x.com/data.csv", csv_data, "text/csv")
        assert page is not None
        assert page.source_type == "csv"
        assert "| Name | Age | City |" in page.content
        assert "| Alice | 30 | Sydney |" in page.content
        assert page.metadata["total_rows"] == 3
        assert page.metadata["columns"] == 3

    def test_extract_tsv(self, extractor):
        tsv_data = b"Col1\tCol2\nA\tB\n"
        page = extractor.extract("https://x.com/data.tsv", tsv_data, "text/plain")
        assert page is not None
        assert "| Col1 | Col2 |" in page.content

    def test_extract_xlsx(self, extractor):
        import io
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Value"])
        ws.append(["Alpha", 100])
        ws.append(["Beta", 200])
        buf = io.BytesIO()
        wb.save(buf)

        page = extractor.extract("https://x.com/data.xlsx", buf.getvalue(), "application/xlsx")
        assert page is not None
        assert page.source_type == "xlsx"
        assert "## Data" in page.content
        assert "| Alpha | 100 |" in page.content
        assert "Data" in page.metadata["sheets"]

    def test_empty_csv_returns_none(self, extractor):
        page = extractor.extract("https://x.com/empty.csv", b"", "text/csv")
        assert page is None


# ─── PptxExtractor ───


class TestPptxExtractor:
    @pytest.fixture
    def extractor(self):
        return PptxExtractor()

    def test_can_handle(self, extractor):
        assert extractor.can_handle("slides.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        assert extractor.can_handle("slides.pptx", "application/octet-stream")
        assert not extractor.can_handle("doc.docx", "application/docx")

    def test_extract_real_pptx(self, extractor):
        from pptx import Presentation
        from pptx.util import Inches
        import io

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Welcome Slide"
        slide.placeholders[1].text = "This is the body content of slide one."

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Second Topic"
        slide2.placeholders[1].text = "Details about the second topic."

        buf = io.BytesIO()
        prs.save(buf)

        page = extractor.extract("https://x.com/talk.pptx", buf.getvalue(), "application/pptx")
        assert page is not None
        assert page.source_type == "pptx"
        assert page.title == "Welcome Slide"
        assert "Second Topic" in page.content
        assert page.metadata["slide_count"] == 2


# ─── RtfExtractor ───


class TestRtfExtractor:
    @pytest.fixture
    def extractor(self):
        return RtfExtractor()

    def test_can_handle(self, extractor):
        assert extractor.can_handle("doc.rtf", "application/rtf")
        assert extractor.can_handle("doc.rtf", "text/rtf")
        assert not extractor.can_handle("doc.pdf", "application/pdf")

    def test_extract_rtf(self, extractor):
        # Simple RTF content
        rtf = rb"{\rtf1\ansi{\fonttbl\f0\fswiss Helvetica;}\f0\pard This is a test RTF document with some content.\par Second paragraph here.\par}"
        page = extractor.extract("https://x.com/doc.rtf", rtf, "application/rtf")
        assert page is not None
        assert page.source_type == "rtf"
        assert "test RTF document" in page.content
        assert page.word_count > 0

    def test_too_short_returns_none(self, extractor):
        rtf = rb"{\rtf1 Hi\par}"
        page = extractor.extract("https://x.com/doc.rtf", rtf, "application/rtf")
        assert page is None


# ─── RssFeedExtractor ───


class TestRssFeedExtractor:
    @pytest.fixture
    def extractor(self):
        return RssFeedExtractor()

    def test_can_handle(self, extractor):
        assert extractor.can_handle("https://x.com/feed", "application/rss+xml")
        assert extractor.can_handle("https://x.com/feed", "application/atom+xml")
        assert extractor.can_handle("https://x.com/feed.rss", "text/xml")
        assert not extractor.can_handle("https://x.com/page", "text/html")

    def test_extract_rss(self, extractor):
        rss = b"""<?xml version="1.0"?>
        <rss version="2.0">
        <channel>
            <title>Test Blog</title>
            <description>A test RSS feed</description>
            <link>https://testblog.com</link>
            <item>
                <title>First Post</title>
                <link>https://testblog.com/first</link>
                <description>This is the first post content.</description>
                <pubDate>Mon, 01 Jan 2024 00:00:00 GMT</pubDate>
                <author>author@test.com</author>
            </item>
            <item>
                <title>Second Post</title>
                <link>https://testblog.com/second</link>
                <description>Second post with more content.</description>
            </item>
        </channel>
        </rss>"""
        page = extractor.extract("https://testblog.com/feed.rss", rss, "application/rss+xml")
        assert page is not None
        assert page.title == "Test Blog"
        assert page.source_type == "rss"
        assert "First Post" in page.content
        assert "Second Post" in page.content
        assert page.metadata["entry_count"] == 2
        assert len(page.links) == 2

    def test_extract_atom(self, extractor):
        atom = b"""<?xml version="1.0" encoding="utf-8"?>
        <feed xmlns="http://www.w3.org/2005/Atom">
            <title>Atom Feed</title>
            <entry>
                <title>Atom Entry</title>
                <link href="https://example.com/entry1"/>
                <summary>Entry summary text.</summary>
                <updated>2024-01-01T00:00:00Z</updated>
                <author><name>AtomAuthor</name></author>
            </entry>
        </feed>"""
        page = extractor.extract("https://x.com/feed", atom, "application/atom+xml")
        assert page is not None
        assert page.title == "Atom Feed"
        assert "Atom Entry" in page.content
        assert "AtomAuthor" in page.content

    def test_empty_feed_returns_none(self, extractor):
        rss = b"""<?xml version="1.0"?><rss version="2.0"><channel></channel></rss>"""
        page = extractor.extract("https://x.com/feed", rss, "application/rss+xml")
        assert page is None


# ─── ScrapedPage metadata ───


class TestScrapedPageMetadata:
    def test_metadata_default_empty(self):
        page = ScrapedPage(
            url="https://example.com", title="Test",
            content="Hello", scraped_at="2024-01-01", word_count=1,
        )
        assert page.metadata == {}

    def test_metadata_preserved(self):
        page = ScrapedPage(
            url="https://example.com", title="Test",
            content="Hello", scraped_at="2024-01-01", word_count=1,
            metadata={"author": "Alice", "page_count": 5},
        )
        assert page.metadata["author"] == "Alice"
        assert page.metadata["page_count"] == 5


# ─── Document extractor registry completeness ───


class TestDocumentExtractorRegistry:
    def test_all_document_extractors_in_registry(self):
        from mini_rag.extractors import _EXTRACTORS
        extractor_types = {type(e).__name__ for e in _EXTRACTORS}
        expected = {
            "ArxivExtractor", "GitHubExtractor", "PDFExtractor",
            "DocxExtractor", "SpreadsheetExtractor", "PptxExtractor",
            "EpubExtractor", "RtfExtractor", "RssFeedExtractor",
            "MarkdownPassthroughExtractor", "GenericExtractor",
        }
        assert expected.issubset(extractor_types)

    def test_generic_is_last(self):
        from mini_rag.extractors import _EXTRACTORS
        assert type(_EXTRACTORS[-1]).__name__ == "GenericExtractor"
