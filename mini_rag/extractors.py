"""
Content extractors for web scraping.

Converts raw HTML and PDF content into clean markdown.
Domain-specific extractors handle known sites (arXiv, GitHub),
with a generic fallback for everything else.
"""

import logging
import re
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Protocol, Tuple
from urllib.parse import urlparse

import requests

logger = logging.getLogger(__name__)

# Optional dependencies — graceful degradation
try:
    from bs4 import BeautifulSoup

    BS4_AVAILABLE = True
except ImportError:
    BeautifulSoup = None
    BS4_AVAILABLE = False

try:
    import fitz  # pymupdf

    PYMUPDF_AVAILABLE = True
except ImportError:
    fitz = None
    PYMUPDF_AVAILABLE = False

try:
    import docx as python_docx

    DOCX_AVAILABLE = True
except ImportError:
    python_docx = None
    DOCX_AVAILABLE = False

try:
    import openpyxl

    OPENPYXL_AVAILABLE = True
except ImportError:
    openpyxl = None
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation

    PPTX_AVAILABLE = True
except ImportError:
    PptxPresentation = None
    PPTX_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub as epub_reader

    EPUB_AVAILABLE = True
except ImportError:
    ebooklib = None
    epub_reader = None
    EPUB_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text

    RTF_AVAILABLE = True
except ImportError:
    rtf_to_text = None
    RTF_AVAILABLE = False

try:
    import feedparser

    FEEDPARSER_AVAILABLE = True
except ImportError:
    feedparser = None
    FEEDPARSER_AVAILABLE = False


@dataclass
class ScrapedPage:
    """A single scraped and extracted page."""

    url: str
    title: str
    content: str  # Clean markdown
    scraped_at: str  # ISO timestamp
    word_count: int
    links: List[str] = field(default_factory=list)  # Outbound links for crawling
    source_type: str = "web"  # "web", "pdf", "arxiv", "github", "docx", etc.
    raw_bytes: Optional[bytes] = None  # Original file (PDF, etc.) for archival
    metadata: Dict = field(default_factory=dict)  # Document metadata (author, pages, etc.)


class ContentExtractor(Protocol):
    """Protocol for content extractors."""

    def can_handle(self, url: str, content_type: str) -> bool: ...

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]: ...


class DirectFetcher(Protocol):
    """Protocol for API-based extractors that fetch + extract in one step.

    These bypass the normal HTTP fetch → extract_content pipeline,
    making their own API calls. Used for sites where:
    - robots.txt blocks scrapers but API access is fine
    - The API returns cleaner data than scraping HTML
    - The site requires JS rendering (YouTube, Reddit)
    """

    def can_handle_url(self, url: str) -> bool: ...

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]: ...


def _slugify(text: str, max_length: int = 60) -> str:
    """Convert text to a filesystem-safe slug."""
    slug = text.lower().strip()
    slug = re.sub(r"[^\w\s-]", "", slug)
    slug = re.sub(r"[\s_]+", "-", slug)
    slug = re.sub(r"-+", "-", slug).strip("-")
    return slug[:max_length]


def _make_frontmatter(page: ScrapedPage) -> str:
    """Generate BOBAI-compatible frontmatter for a scraped page."""
    lines = [
        "---",
        "profile: scraped",
        'generator: "fss-mini-rag-scraper"',
        f'title: "{page.title}"',
        f'source_url: "{page.url}"',
        f'scraped_at: "{page.scraped_at}"',
        f"word_count: {page.word_count}",
        f'source_type: "{page.source_type}"',
        "content_quality: 1.0",
    ]
    # Include document metadata if present
    for key, val in page.metadata.items():
        if val is not None and val != "":
            # Sanitize value for YAML
            safe_val = str(val).replace('"', '\\"')
            lines.append(f'{key}: "{safe_val}"')
    lines.append("---\n\n")
    return "\n".join(lines)


def save_scraped_page(page: ScrapedPage, output_dir: Path) -> Path:
    """Save a scraped page as markdown with frontmatter.

    Returns the path to the saved file.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(page.title) or "untitled"
    filename = f"{slug}.md"
    filepath = output_dir / filename

    # Avoid collisions
    counter = 1
    while filepath.exists():
        filename = f"{slug}-{counter}.md"
        filepath = output_dir / filename
        counter += 1

    frontmatter = _make_frontmatter(page)
    footer = (
        f"\n\n---\n*Source: [{page.url}]({page.url}) "
        f"— scraped {page.scraped_at[:10]}*\n"
    )

    content = frontmatter + f"# {page.title}\n\n" + page.content + footer
    filepath.write_text(content, encoding="utf-8")
    logger.info(f"Saved: {filepath}")

    # Save original binary (PDF, etc.) alongside the markdown
    if page.raw_bytes and page.source_type in ("pdf", "arxiv"):
        originals_dir = output_dir / "originals"
        originals_dir.mkdir(exist_ok=True)
        ext = ".pdf"  # Default for pdf/arxiv
        original_path = originals_dir / f"{slug}{ext}"
        if not original_path.exists():
            original_path.write_bytes(page.raw_bytes)
            logger.info(f"Saved original: {original_path}")

    return filepath


class GenericExtractor:
    """Generic HTML to markdown extractor using BeautifulSoup.

    Strips navigation, footers, scripts, and other noise.
    Extracts main content and converts to clean markdown.
    """

    # Elements to remove entirely
    REMOVE_TAGS = {
        "script", "style", "nav", "footer", "header", "aside",
        "iframe", "noscript", "svg", "form", "button",
    }

    # Elements likely to contain main content
    CONTENT_SELECTORS = [
        "article",
        "main",
        '[role="main"]',
        ".post-content",
        ".article-content",
        ".entry-content",
        "#content",
        ".content",
    ]

    def can_handle(self, url: str, content_type: str) -> bool:
        return "text/html" in content_type

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not BS4_AVAILABLE:
            logger.warning(
                "beautifulsoup4 not installed. Install with: pip install beautifulsoup4"
            )
            return None

        try:
            # Detect encoding from content-type or let bs4 handle it
            soup = BeautifulSoup(raw.decode("utf-8", errors="replace"), "html.parser")

            # Extract title
            title = self._extract_title(soup, url)

            # Remove noise elements
            for tag_name in self.REMOVE_TAGS:
                for tag in soup.find_all(tag_name):
                    tag.decompose()

            # Find main content area
            content_element = self._find_content(soup)

            # Convert to markdown
            markdown = self._html_to_markdown(content_element)

            if not markdown or len(markdown.strip()) < 50:
                logger.debug(f"Insufficient content from {url}")
                return None

            # Extract links for crawling
            links = self._extract_links(soup, url)

            word_count = len(markdown.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=markdown,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=links,
                source_type="web",
            )

        except Exception as e:
            logger.error(f"Extraction failed for {url}: {e}")
            return None

    def _extract_title(self, soup, url: str) -> str:
        """Extract page title from various sources."""
        # Try og:title first (usually cleanest)
        og_title = soup.find("meta", property="og:title")
        if og_title and og_title.get("content"):
            return og_title["content"].strip()

        # Try <title> tag
        if soup.title and soup.title.string:
            title = soup.title.string.strip()
            # Remove common suffixes like " | Site Name" or " - Site Name"
            title = re.split(r"\s*[\|\-–—]\s*", title)[0].strip()
            if title:
                return title

        # Try first h1
        h1 = soup.find("h1")
        if h1:
            return h1.get_text(strip=True)

        # Fallback to domain
        return urlparse(url).netloc

    def _find_content(self, soup):
        """Find the main content element."""
        # Try known content selectors
        for selector in self.CONTENT_SELECTORS:
            element = soup.select_one(selector)
            if element and len(element.get_text(strip=True)) > 100:
                return element

        # Fallback: find the div with most text content
        best = None
        best_length = 0
        for div in soup.find_all(["div", "section"]):
            text = div.get_text(strip=True)
            if len(text) > best_length:
                best_length = len(text)
                best = div

        return best or soup.body or soup

    def _html_to_markdown(self, element) -> str:
        """Convert an HTML element to clean markdown."""
        if element is None:
            return ""

        lines = []
        self._process_element(element, lines)

        # Clean up excessive blank lines
        text = "\n".join(lines)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _process_element(self, element, lines: list, depth: int = 0):
        """Recursively process HTML elements into markdown lines."""
        if hasattr(element, "name"):
            tag = element.name

            if tag in self.REMOVE_TAGS:
                return

            # Headings
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(tag[1])
                text = element.get_text(strip=True)
                if text:
                    lines.append("")
                    lines.append(f"{'#' * level} {text}")
                    lines.append("")
                return

            # Paragraphs
            if tag == "p":
                text = element.get_text(strip=True)
                if text:
                    lines.append("")
                    lines.append(text)
                return

            # Lists
            if tag in ("ul", "ol"):
                lines.append("")
                for i, li in enumerate(element.find_all("li", recursive=False)):
                    prefix = f"{i + 1}." if tag == "ol" else "-"
                    text = li.get_text(strip=True)
                    if text:
                        lines.append(f"{prefix} {text}")
                lines.append("")
                return

            # Code blocks
            if tag == "pre":
                code = element.find("code")
                text = (code or element).get_text()
                lang = ""
                if code and code.get("class"):
                    for cls in code["class"]:
                        if cls.startswith("language-"):
                            lang = cls[9:]
                            break
                lines.append("")
                lines.append(f"```{lang}")
                lines.append(text.rstrip())
                lines.append("```")
                lines.append("")
                return

            # Inline code
            if tag == "code" and element.parent and element.parent.name != "pre":
                text = element.get_text()
                if text:
                    lines.append(f"`{text}`")
                return

            # Blockquotes
            if tag == "blockquote":
                text = element.get_text(strip=True)
                if text:
                    lines.append("")
                    for line in text.split("\n"):
                        lines.append(f"> {line.strip()}")
                    lines.append("")
                return

            # Links
            if tag == "a":
                href = element.get("href", "")
                text = element.get_text(strip=True)
                if text and href and not href.startswith("#"):
                    lines.append(f"[{text}]({href})")
                elif text:
                    lines.append(text)
                return

            # Images
            if tag == "img":
                alt = element.get("alt", "")
                src = element.get("src", "")
                if src:
                    lines.append(f"![{alt}]({src})")
                return

            # Tables
            if tag == "table":
                self._process_table(element, lines)
                return

            # For other block elements, recurse into children
            if tag in ("div", "section", "article", "main", "span", "figure",
                        "figcaption", "details", "summary"):
                for child in element.children:
                    self._process_element(child, lines, depth + 1)
                return

        # NavigableString — raw text
        elif hasattr(element, "string") and element.string:
            text = element.string.strip()
            if text:
                lines.append(text)

    def _process_table(self, table, lines: list):
        """Convert an HTML table to markdown."""
        rows = table.find_all("tr")
        if not rows:
            return

        lines.append("")
        for i, row in enumerate(rows):
            cells = row.find_all(["th", "td"])
            cell_texts = [c.get_text(strip=True).replace("|", "\\|") for c in cells]
            lines.append("| " + " | ".join(cell_texts) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * len(cell_texts)) + " |")
        lines.append("")

    def _extract_links(self, soup, base_url: str) -> List[str]:
        """Extract outbound links from the page."""
        from urllib.parse import urljoin

        links = []
        seen = set()
        for a_tag in soup.find_all("a", href=True):
            href = a_tag["href"]
            if href.startswith(("#", "javascript:", "mailto:")):
                continue
            full_url = urljoin(base_url, href)
            if full_url not in seen:
                seen.add(full_url)
                links.append(full_url)

        return links


class PDFExtractor:
    """Extract text from PDF files using pymupdf (fitz).

    Uses font-size analysis to detect headings and document structure.
    Converts to clean markdown with proper heading hierarchy.
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        return (
            "application/pdf" in content_type
            or url.lower().endswith(".pdf")
        )

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not PYMUPDF_AVAILABLE:
            logger.warning(
                "pymupdf not installed. Install with: pip install pymupdf"
            )
            return None

        try:
            doc = fitz.open(stream=raw, filetype="pdf")
            metadata = doc.metadata or {}
            title = metadata.get("title", "").strip()
            page_count = len(doc)

            markdown, stats = self._extract_structured(doc)
            doc.close()

            if not markdown:
                logger.debug(f"No text extracted from PDF: {url}")
                return None

            markdown = self._clean_pdf_text(markdown)

            if not title:
                first_line = markdown.split("\n")[0].strip()
                if first_line and len(first_line) < 200:
                    title = first_line.lstrip("# ").strip()
                else:
                    title = Path(urlparse(url).path).stem or "Untitled PDF"

            word_count = len(markdown.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=markdown,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="pdf",
                raw_bytes=raw,
            )

        except Exception as e:
            logger.error(f"PDF extraction failed for {url}: {e}")
            return None

    def _extract_structured(self, doc) -> Tuple[str, Dict]:
        """Extract text with structure detection from font metadata.

        Two-pass approach:
        1. Build font-size histogram to identify body vs heading sizes
        2. Extract text, converting larger fonts to markdown headings
        """
        # Phase 1: font-size histogram
        size_counter: Counter = Counter()
        for page in doc:
            for block in page.get_text("dict")["blocks"]:
                if "lines" not in block:
                    continue
                for line in block["lines"]:
                    for span in line["spans"]:
                        text = span["text"].strip()
                        if text:
                            size_counter[round(span["size"], 1)] += len(text)

        if not size_counter:
            return "", {}

        body_size = size_counter.most_common(1)[0][0]

        # Heading sizes: anything > body_size * 1.15, ranked largest first
        heading_sizes = sorted(
            [sz for sz in size_counter if sz > body_size * 1.15],
            reverse=True,
        )
        heading_map: Dict[float, int] = {}
        for i, sz in enumerate(heading_sizes[:3]):
            heading_map[sz] = i + 1  # 1=H2, 2=H3, 3=H4 (offset +1 in output)

        # Phase 2: extract with structure
        pages_md = []
        for page in doc:
            lines_out: List[str] = []
            prev_was_heading = False

            for block in page.get_text("dict")["blocks"]:
                if "lines" not in block:
                    continue
                for line in block["lines"]:
                    line_text = ""
                    line_size = None
                    is_bold = False

                    for span in line["spans"]:
                        line_text += span["text"]
                        if span["text"].strip():
                            line_size = round(span["size"], 1)
                            is_bold = bool(span["flags"] & 16)

                    line_text = line_text.rstrip()
                    stripped = line_text.strip()
                    if not stripped:
                        continue

                    # Heading by font size
                    if line_size in heading_map:
                        level = heading_map[line_size] + 1  # +1 so title=##, sections=###
                        lines_out.append("")
                        lines_out.append(f"{'#' * level} {stripped}")
                        lines_out.append("")
                        prev_was_heading = True
                        continue

                    # Bold subheading (same body size but bold, short, starts a block)
                    if is_bold and line_size == body_size and len(stripped) < 100:
                        # Only treat as subheading if it looks like a title (not mid-paragraph bold)
                        if prev_was_heading or not lines_out or not lines_out[-1].strip():
                            lines_out.append("")
                            lines_out.append(f"### {stripped}")
                            lines_out.append("")
                            prev_was_heading = True
                            continue

                    lines_out.append(line_text)
                    prev_was_heading = False

            page_text = "\n".join(lines_out)
            if page_text.strip():
                pages_md.append(page_text)

        stats = {"body_size": body_size, "heading_sizes": heading_sizes}
        return "\n\n***\n\n".join(pages_md), stats

    def _clean_pdf_text(self, text: str) -> str:
        """Clean up extracted PDF text.

        Handles: paragraph reconstruction, dehyphenation, ALL-CAPS section
        detection, page numbers, excessive whitespace.
        """
        # Dehyphenation: fix word-break hyphens at line ends
        text = re.sub(r"([a-z])-\n([a-z])", r"\1\2", text)

        # Reconstruct paragraphs from line-per-block PDF output
        text = self._join_paragraphs(text)

        lines = text.split("\n")
        cleaned = []
        for line in lines:
            line = line.rstrip()
            line = re.sub(r" {2,}", " ", line)  # Collapse multiple spaces
            stripped = line.strip()

            # ALL-CAPS section detection (academic papers)
            # Must be standalone line, 5-80 chars, actual words, not already a heading
            if (
                stripped
                and not stripped.startswith("#")
                and 5 <= len(stripped) <= 80
                and re.match(r"^[A-Z][A-Z\s,&:;\'\-]{3,}$", stripped)
                and sum(1 for c in stripped if c.isalpha()) > len(stripped) * 0.5
            ):
                cleaned.append("")
                cleaned.append(f"## {stripped}")
                cleaned.append("")
                continue

            cleaned.append(line)

        text = "\n".join(cleaned)

        # Remove standalone page numbers
        text = re.sub(r"\n\s*\d{1,4}\s*\n", "\n", text)

        # Collapse excessive blank lines
        text = re.sub(r"\n{4,}", "\n\n***\n\n", text)
        text = re.sub(r"\n{3,}", "\n\n", text)

        return text.strip()

    @staticmethod
    def _join_paragraphs(text: str) -> str:
        """Reconstruct paragraphs from line-per-block PDF output.

        PDF extraction produces one line per text block. Consecutive body text
        lines that belong to the same paragraph are joined with spaces.
        Paragraph breaks are inserted at sentence boundaries when the
        accumulated paragraph is long enough to be a complete thought.
        """
        lines = text.split("\n")
        result: List[str] = []
        current_para: List[str] = []
        current_len = 0

        def _flush_para():
            nonlocal current_len
            if current_para:
                result.append(" ".join(current_para))
                result.append("")  # blank line = paragraph break
                current_para.clear()
                current_len = 0

        for line in lines:
            stripped = line.strip()

            # Empty line — flush current paragraph
            if not stripped:
                _flush_para()
                continue

            # Headings, separators — always their own block with surrounding breaks
            if stripped.startswith("#") or stripped == "***":
                _flush_para()
                result.append(line)
                result.append("")
                continue

            # If current paragraph is empty, start a new one
            if not current_para:
                current_para.append(stripped)
                current_len = len(stripped)
                continue

            # Decide: join to current paragraph or start new one?
            # Look at the accumulated paragraph so far
            joined_so_far = " ".join(current_para)
            last_char = joined_so_far.rstrip()[-1] if joined_so_far.rstrip() else ""

            # Sentence-end detection: flush paragraph when ALL conditions met:
            # 1. Current para ends with sentence-ending punctuation
            # 2. Next line starts with uppercase (new sentence)
            # 3. Current para is at least 80 chars (not just a short fragment)
            is_sentence_end = last_char in '.!?"\')\u201d'
            starts_new = stripped[0].isupper() or stripped[0] in "\"\u2018\u201c(-\u2022\u2013\u2014"
            long_enough = current_len > 80

            if is_sentence_end and starts_new and long_enough:
                _flush_para()
                current_para.append(stripped)
                current_len = len(stripped)
                continue

            # Otherwise: continuation — join to current paragraph
            current_para.append(stripped)
            current_len += len(stripped) + 1

        _flush_para()
        return "\n".join(result)


class ArxivExtractor:
    """Extract content from arXiv pages.

    Handles:
    - Abstract pages (arxiv.org/abs/...) — extracts title, authors, abstract, metadata
    - PDF links (arxiv.org/pdf/...) — delegates to PDFExtractor
    """

    _pdf_extractor = PDFExtractor()

    # URL patterns we can handle
    _HANDLED_PREFIXES = ("/abs/", "/pdf/", "/html/")

    def can_handle(self, url: str, content_type: str) -> bool:
        parsed = urlparse(url)
        if "arxiv.org" not in parsed.netloc:
            return False
        # Only handle article pages, not listings or search
        return any(parsed.path.startswith(p) for p in self._HANDLED_PREFIXES)

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        # PDF on arXiv — delegate
        if self._pdf_extractor.can_handle(url, content_type):
            page = self._pdf_extractor.extract(url, raw, content_type)
            if page:
                page.source_type = "arxiv"
            return page

        if not BS4_AVAILABLE:
            return None

        # /html/ pages — full article HTML, use generic extraction with arxiv metadata
        parsed_path = urlparse(url).path
        if parsed_path.startswith("/html/"):
            return self._extract_html_article(url, raw)

        try:
            soup = BeautifulSoup(raw.decode("utf-8", errors="replace"), "html.parser")

            # Title
            title_el = soup.select_one("h1.title")
            title = title_el.get_text(strip=True).replace("Title:", "").strip() if title_el else ""

            # Authors
            authors_el = soup.select_one(".authors")
            authors = authors_el.get_text(strip=True).replace("Authors:", "").strip() if authors_el else ""

            # Abstract
            abstract_el = soup.select_one("blockquote.abstract")
            abstract = abstract_el.get_text(strip=True).replace("Abstract:", "").strip() if abstract_el else ""

            # Subjects/categories
            subjects_el = soup.select_one(".subjects")
            subjects = subjects_el.get_text(strip=True) if subjects_el else ""

            # Submission date
            date_el = soup.select_one(".dateline")
            dateline = date_el.get_text(strip=True) if date_el else ""

            # Build markdown
            parts = [f"**Authors:** {authors}"] if authors else []
            if dateline:
                parts.append(f"**Submitted:** {dateline}")
            if subjects:
                parts.append(f"**Subjects:** {subjects}")
            parts.append("")
            parts.append("## Abstract")
            parts.append("")
            parts.append(abstract)

            # Get PDF link
            pdf_link = soup.select_one('a[href*="/pdf/"]')
            if pdf_link:
                pdf_url = pdf_link.get("href", "")
                if pdf_url.startswith("/"):
                    pdf_url = f"https://arxiv.org{pdf_url}"
                parts.append("")
                parts.append(f"**PDF:** [{pdf_url}]({pdf_url})")

            content = "\n".join(parts)
            word_count = len(content.split())

            if not title and not abstract:
                return None

            return ScrapedPage(
                url=url,
                title=title or "Untitled arXiv Paper",
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="arxiv",
            )

        except Exception as e:
            logger.error(f"arXiv extraction failed for {url}: {e}")
            return None

    def _extract_html_article(self, url: str, raw: bytes) -> Optional[ScrapedPage]:
        """Extract from arXiv /html/ full-text article pages."""
        try:
            soup = BeautifulSoup(raw.decode("utf-8", errors="replace"), "html.parser")

            # Title from the page
            title_el = soup.select_one("h1.ltx_title") or soup.select_one("title")
            title = title_el.get_text(strip=True) if title_el else ""

            # Authors
            authors_el = soup.select_one(".ltx_authors")
            authors = authors_el.get_text(strip=True) if authors_el else ""

            # Abstract
            abstract_el = soup.select_one(".ltx_abstract")
            abstract = ""
            if abstract_el:
                abstract = abstract_el.get_text(strip=True)
                abstract = abstract.replace("Abstract.", "").replace("Abstract", "").strip()

            # Remove nav/header elements
            for tag in soup.select("nav, header, footer, .ltx_page_footer, .ltx_page_header"):
                tag.decompose()

            # Extract the main body
            extractor = GenericExtractor()
            body = soup.select_one(".ltx_page_content") or soup.select_one("article") or soup.body
            body_md = extractor._html_to_markdown(body) if body else ""

            parts = []
            if authors:
                parts.append(f"**Authors:** {authors}")
            if abstract:
                parts.append(f"\n## Abstract\n\n{abstract}")
            if body_md:
                parts.append(f"\n{body_md}")

            content = "\n".join(parts)
            word_count = len(content.split())

            if not title and not abstract:
                return None

            return ScrapedPage(
                url=url,
                title=title or "Untitled arXiv Paper",
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="arxiv",
            )

        except Exception as e:
            logger.error(f"arXiv HTML extraction failed for {url}: {e}")
            return None


class GitHubExtractor:
    """Extract content from GitHub pages.

    Handles:
    - Repository pages (github.com/user/repo) — extracts README
    - File views — extracts code content
    - Focuses on the actual content, stripping GitHub chrome
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        """Only handle GitHub repo pages (user/repo), not listing pages."""
        if "github.com" not in urlparse(url).netloc or "text/html" not in content_type:
            return False
        path_parts = [p for p in urlparse(url).path.split("/") if p]
        # Need at least user/repo — skip /, /trending, /explore, /features etc.
        return len(path_parts) >= 2

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not BS4_AVAILABLE:
            return None

        try:
            soup = BeautifulSoup(raw.decode("utf-8", errors="replace"), "html.parser")

            parsed_url = urlparse(url)
            path_parts = [p for p in parsed_url.path.split("/") if p]

            # Repository main page — extract README
            readme_el = soup.select_one("article.markdown-body")
            if readme_el:
                title = self._extract_repo_title(soup, path_parts)
                content = self._readme_to_markdown(readme_el)

                # Extract repo metadata
                description_el = soup.select_one("p.f4.my-3")
                description = description_el.get_text(strip=True) if description_el else ""

                # Topics/tags
                topics = [t.get_text(strip=True) for t in soup.select("a.topic-tag")]

                parts = []
                if description:
                    parts.append(f"*{description}*\n")
                if topics:
                    parts.append(f"**Topics:** {', '.join(topics)}\n")
                parts.append(content)

                full_content = "\n".join(parts)
                word_count = len(full_content.split())

                return ScrapedPage(
                    url=url,
                    title=title,
                    content=full_content,
                    scraped_at=datetime.now().isoformat(),
                    word_count=word_count,
                    links=self._extract_links(soup, url),
                    source_type="github",
                )

            # Code file view
            code_el = soup.select_one("table.highlight")
            if code_el:
                title = path_parts[-1] if path_parts else "GitHub File"
                code_text = code_el.get_text()
                repo_name = "/".join(path_parts[:2]) if len(path_parts) >= 2 else ""

                content = f"**Repository:** {repo_name}\n\n```\n{code_text}\n```"
                word_count = len(content.split())

                return ScrapedPage(
                    url=url,
                    title=f"{repo_name}/{title}",
                    content=content,
                    scraped_at=datetime.now().isoformat(),
                    word_count=word_count,
                    links=[],
                    source_type="github",
                )

            return None

        except Exception as e:
            logger.error(f"GitHub extraction failed for {url}: {e}")
            return None

    def _extract_repo_title(self, soup, path_parts: list) -> str:
        """Extract repository name from GitHub page."""
        if len(path_parts) >= 2:
            return f"{path_parts[0]}/{path_parts[1]}"
        title_el = soup.select_one("strong[itemprop='name'] a")
        if title_el:
            return title_el.get_text(strip=True)
        return "GitHub Repository"

    def _readme_to_markdown(self, article) -> str:
        """Convert GitHub's rendered README back to clean markdown."""
        # GitHub's markdown-body is already rendered HTML — use generic conversion
        extractor = GenericExtractor()
        return extractor._html_to_markdown(article)

    def _extract_links(self, soup, base_url: str) -> List[str]:
        """Extract relevant links from GitHub page."""
        from urllib.parse import urljoin
        links = []
        for a in soup.select("article.markdown-body a[href]"):
            href = a["href"]
            if href.startswith(("#", "javascript:")):
                continue
            links.append(urljoin(base_url, href))
        return links


class MarkdownPassthroughExtractor:
    """Pass through markdown and plain text content with minimal processing.

    Handles Content-Type: text/markdown, text/plain, and URLs ending
    in .md or .txt that would otherwise get no extractor.
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        if any(t in ct for t in ("text/markdown", "text/x-markdown")):
            return True
        if "text/plain" in ct:
            return True
        path = urlparse(url).path.lower()
        return path.endswith(".md") or path.endswith(".txt")

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        try:
            text = raw.decode("utf-8", errors="replace")
        except Exception:
            return None

        if not text or len(text.strip()) < 20:
            return None

        # Extract title: first # heading, or first non-empty line, or filename
        title = None
        for line in text.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("#"):
                title = stripped.lstrip("#").strip()
                break
            if len(stripped) < 200:
                title = stripped
            break

        if not title:
            path = urlparse(url).path
            title = Path(path).stem or urlparse(url).netloc

        word_count = len(text.split())

        return ScrapedPage(
            url=url,
            title=title,
            content=text.strip(),
            scraped_at=datetime.now().isoformat(),
            word_count=word_count,
            links=[],
            source_type="markdown",
        )


# ─── Document extractors ───


class DocxExtractor:
    """Extract content from DOCX files using python-docx.

    Extracts headings, paragraphs, tables, and document metadata
    (author, created date, last modified, company).
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        if "officedocument.wordprocessingml" in ct:
            return True
        if "application/msword" in ct:
            return True
        return urlparse(url).path.lower().endswith(".docx")

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not installed. Install with: pip install python-docx")
            return None

        try:
            import io
            doc = python_docx.Document(io.BytesIO(raw))

            # Extract metadata
            props = doc.core_properties
            meta = {}
            if props.author:
                meta["author"] = props.author
            if props.created:
                meta["created"] = props.created.isoformat()
            if props.modified:
                meta["modified"] = props.modified.isoformat()
            if props.last_modified_by:
                meta["last_modified_by"] = props.last_modified_by
            if props.category:
                meta["category"] = props.category
            if props.subject:
                meta["subject"] = props.subject

            title = props.title or ""
            lines = []

            for element in doc.element.body:
                tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

                if tag == "p":
                    para = None
                    for p in doc.paragraphs:
                        if p._element is element:
                            para = p
                            break
                    if para is None:
                        continue

                    style = para.style.name if para.style else ""
                    text = para.text.strip()
                    if not text:
                        continue

                    if "Heading 1" in style:
                        if not title:
                            title = text
                        lines.append(f"\n## {text}\n")
                    elif "Heading 2" in style:
                        lines.append(f"\n### {text}\n")
                    elif "Heading 3" in style:
                        lines.append(f"\n#### {text}\n")
                    elif "List" in style:
                        lines.append(f"- {text}")
                    else:
                        lines.append(text)

                elif tag == "tbl":
                    for table in doc.tables:
                        if table._element is element:
                            lines.append(self._table_to_markdown(table))
                            break

            if not title:
                title = Path(urlparse(url).path).stem or "Untitled Document"

            content = "\n".join(lines).strip()
            word_count = len(content.split())

            meta["page_count"] = len(doc.sections)

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="docx",
                raw_bytes=raw,
                metadata=meta,
            )

        except Exception as e:
            logger.error(f"DOCX extraction failed for {url}: {e}")
            return None

    @staticmethod
    def _table_to_markdown(table) -> str:
        """Convert a python-docx table to markdown."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if len(rows) >= 1:
            # Insert header separator after first row
            col_count = len(table.rows[0].cells)
            rows.insert(1, "| " + " | ".join(["---"] * col_count) + " |")
        return "\n" + "\n".join(rows) + "\n"


class SpreadsheetExtractor:
    """Extract content from XLSX and CSV files.

    XLSX: uses openpyxl to read sheets and convert to markdown tables.
    CSV: uses stdlib csv module.
    Tracks sheet names, row/column counts in metadata.
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        path = urlparse(url).path.lower()
        if "spreadsheetml" in ct or "text/csv" in ct or "application/csv" in ct:
            return True
        return path.endswith((".xlsx", ".csv", ".tsv"))

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        path = urlparse(url).path.lower()
        ct = content_type.lower()

        if path.endswith(".csv") or path.endswith(".tsv") or "csv" in ct:
            return self._extract_csv(url, raw, path.endswith(".tsv"))

        if not OPENPYXL_AVAILABLE:
            logger.warning("openpyxl not installed. Install with: pip install openpyxl")
            return None

        return self._extract_xlsx(url, raw)

    def _extract_csv(self, url: str, raw: bytes, is_tsv: bool = False) -> Optional[ScrapedPage]:
        """Extract CSV/TSV to markdown table."""
        import csv
        import io

        try:
            text = raw.decode("utf-8", errors="replace")
            delimiter = "\t" if is_tsv else ","
            reader = csv.reader(io.StringIO(text), delimiter=delimiter)
            rows = list(reader)

            if not rows:
                return None

            # Limit to first 200 rows for sanity
            display_rows = rows[:200]
            lines = []
            for i, row in enumerate(display_rows):
                cells = [c.strip().replace("|", "\\|") for c in row]
                lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    lines.append("| " + " | ".join(["---"] * len(cells)) + " |")

            content = "\n".join(lines)
            word_count = len(content.split())
            title = Path(urlparse(url).path).stem or "Spreadsheet"
            meta = {"total_rows": len(rows), "columns": len(rows[0]) if rows else 0}
            if len(rows) > 200:
                meta["truncated_at"] = 200

            return ScrapedPage(
                url=url, title=title, content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count, links=[],
                source_type="csv", metadata=meta,
            )
        except Exception as e:
            logger.error(f"CSV extraction failed for {url}: {e}")
            return None

    def _extract_xlsx(self, url: str, raw: bytes) -> Optional[ScrapedPage]:
        """Extract XLSX sheets to markdown tables."""
        import io

        try:
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts = []
            meta = {"sheets": [], "total_rows": 0}

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(max_row=200, values_only=True):
                    cells = [str(c).strip().replace("|", "\\|") if c is not None else "" for c in row]
                    if any(cells):
                        rows.append(cells)

                if not rows:
                    continue

                meta["sheets"].append(sheet_name)
                meta["total_rows"] += ws.max_row or 0

                lines = [f"## {sheet_name}\n"]
                for i, cells in enumerate(rows):
                    lines.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                parts.append("\n".join(lines))

            wb.close()

            if not parts:
                return None

            content = "\n\n".join(parts)
            word_count = len(content.split())
            title = Path(urlparse(url).path).stem or "Spreadsheet"

            return ScrapedPage(
                url=url, title=title, content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count, links=[],
                source_type="xlsx", metadata=meta,
            )
        except Exception as e:
            logger.error(f"XLSX extraction failed for {url}: {e}")
            return None


class PptxExtractor:
    """Extract content from PowerPoint PPTX files.

    Extracts slide titles, text content, speaker notes, and tables.
    Tracks slide count and layout info in metadata.
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        if "presentationml" in ct:
            return True
        return urlparse(url).path.lower().endswith(".pptx")

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            return None

        try:
            import io
            prs = PptxPresentation(io.BytesIO(raw))

            parts = []
            slide_count = 0
            title = ""

            for slide in prs.slides:
                slide_count += 1
                slide_title = ""
                slide_text = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if not text:
                                continue
                            if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                                slide_title = text
                            else:
                                slide_text.append(text)

                    if shape.has_table:
                        table_lines = []
                        for i, row in enumerate(shape.table.rows):
                            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                            table_lines.append("| " + " | ".join(cells) + " |")
                            if i == 0:
                                table_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        slide_text.append("\n".join(table_lines))

                # Speaker notes
                notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()

                if not slide_title and not slide_text:
                    continue

                if not title and slide_title:
                    title = slide_title

                slide_md = f"## Slide {slide_count}"
                if slide_title:
                    slide_md += f": {slide_title}"
                slide_md += "\n\n"
                if slide_text:
                    slide_md += "\n".join(slide_text)
                if notes:
                    slide_md += f"\n\n*Speaker notes: {notes}*"

                parts.append(slide_md)

            if not parts:
                return None

            if not title:
                title = Path(urlparse(url).path).stem or "Presentation"

            content = "\n\n---\n\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url, title=title, content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count, links=[],
                source_type="pptx",
                raw_bytes=raw,
                metadata={"slide_count": slide_count},
            )

        except Exception as e:
            logger.error(f"PPTX extraction failed for {url}: {e}")
            return None


class EpubExtractor:
    """Extract content from EPUB ebook files.

    Extracts chapters as sections, book metadata (author, publisher,
    language, date). Converts XHTML content to markdown.
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        if "application/epub" in ct:
            return True
        return urlparse(url).path.lower().endswith(".epub")

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not EPUB_AVAILABLE:
            logger.warning("ebooklib not installed. Install with: pip install ebooklib")
            return None

        try:
            import io
            book = epub_reader.read_epub(io.BytesIO(raw))

            # Extract metadata
            meta = {}
            title = ""
            for m in book.get_metadata("DC", "title"):
                title = m[0]
                break
            for m in book.get_metadata("DC", "creator"):
                meta["author"] = m[0]
                break
            for m in book.get_metadata("DC", "publisher"):
                meta["publisher"] = m[0]
                break
            for m in book.get_metadata("DC", "language"):
                meta["language"] = m[0]
                break
            for m in book.get_metadata("DC", "date"):
                meta["date"] = m[0]
                break

            # Extract chapter content
            parts = []
            chapter_count = 0
            extractor = GenericExtractor()

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    chapter_count += 1
                    content_bytes = item.get_content()
                    if not content_bytes:
                        continue

                    if BS4_AVAILABLE:
                        soup = BeautifulSoup(content_bytes.decode("utf-8", errors="replace"), "html.parser")
                        # Remove style/script
                        for tag in soup.find_all(["style", "script"]):
                            tag.decompose()
                        md = extractor._html_to_markdown(soup.body or soup)
                        if md and len(md.strip()) > 20:
                            parts.append(md)

            if not parts:
                return None

            if not title:
                title = Path(urlparse(url).path).stem or "Ebook"

            meta["chapter_count"] = chapter_count
            content = "\n\n---\n\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url, title=title, content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count, links=[],
                source_type="epub",
                raw_bytes=raw,
                metadata=meta,
            )

        except Exception as e:
            logger.error(f"EPUB extraction failed for {url}: {e}")
            return None


class RtfExtractor:
    """Extract content from RTF files using striprtf."""

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        if "application/rtf" in ct or "text/rtf" in ct:
            return True
        return urlparse(url).path.lower().endswith(".rtf")

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not RTF_AVAILABLE:
            logger.warning("striprtf not installed. Install with: pip install striprtf")
            return None

        try:
            text = raw.decode("utf-8", errors="replace")
            content = rtf_to_text(text)

            if not content or len(content.strip()) < 20:
                return None

            # Extract title from first line
            title = None
            for line in content.split("\n"):
                stripped = line.strip()
                if stripped and len(stripped) < 200:
                    title = stripped
                    break
            if not title:
                title = Path(urlparse(url).path).stem or "RTF Document"

            word_count = len(content.split())

            return ScrapedPage(
                url=url, title=title, content=content.strip(),
                scraped_at=datetime.now().isoformat(),
                word_count=word_count, links=[],
                source_type="rtf",
            )

        except Exception as e:
            logger.error(f"RTF extraction failed for {url}: {e}")
            return None


class RssFeedExtractor:
    """Extract content from RSS and Atom feeds.

    Parses feed entries, extracting titles, content, dates, authors.
    Returns structured markdown with entry summaries and links.
    """

    def can_handle(self, url: str, content_type: str) -> bool:
        ct = content_type.lower()
        if any(t in ct for t in ("rss", "atom", "xml")):
            # Check if it's actually a feed (not just generic XML)
            if "rss" in ct or "atom" in ct or "feed" in ct:
                return True
        path = urlparse(url).path.lower()
        return path.endswith((".rss", ".atom", ".feed"))

    def extract(self, url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
        if not FEEDPARSER_AVAILABLE:
            logger.warning("feedparser not installed. Install with: pip install feedparser")
            return None

        try:
            feed = feedparser.parse(raw)

            if not feed.entries and not feed.feed.get("title"):
                return None

            feed_title = feed.feed.get("title", "")
            feed_desc = feed.feed.get("description", "")
            feed_link = feed.feed.get("link", "")

            parts = []
            if feed_desc:
                parts.append(f"*{feed_desc}*\n")
            if feed_link:
                parts.append(f"**Site:** [{feed_link}]({feed_link})\n")

            meta = {
                "feed_type": feed.version or "unknown",
                "entry_count": len(feed.entries),
            }

            for entry in feed.entries[:50]:  # Cap at 50 entries
                entry_title = entry.get("title", "Untitled")
                entry_link = entry.get("link", "")
                entry_date = entry.get("published", entry.get("updated", ""))
                entry_author = entry.get("author", "")

                # Get content or summary
                entry_content = ""
                if entry.get("content"):
                    entry_content = entry.content[0].get("value", "")
                elif entry.get("summary"):
                    entry_content = entry.summary

                # Convert HTML content to text
                if entry_content and BS4_AVAILABLE and "<" in entry_content:
                    soup = BeautifulSoup(entry_content, "html.parser")
                    entry_content = soup.get_text(strip=True)

                # Truncate long entries
                if len(entry_content) > 500:
                    entry_content = entry_content[:500] + "..."

                header = f"### {entry_title}"
                if entry_link:
                    header = f"### [{entry_title}]({entry_link})"

                entry_parts = [header]
                meta_line = []
                if entry_date:
                    meta_line.append(entry_date)
                if entry_author:
                    meta_line.append(f"by {entry_author}")
                if meta_line:
                    entry_parts.append(f"*{' | '.join(meta_line)}*")
                if entry_content:
                    entry_parts.append(f"\n{entry_content}")

                parts.append("\n".join(entry_parts))

            if not parts:
                return None

            title = feed_title or Path(urlparse(url).path).stem or "Feed"
            content = "\n\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url, title=title, content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[e.get("link", "") for e in feed.entries[:50] if e.get("link")],
                source_type="rss",
                metadata=meta,
            )

        except Exception as e:
            logger.error(f"RSS/Atom extraction failed for {url}: {e}")
            return None


# ─── API-based direct fetchers ───
# These bypass the normal fetch → extract pipeline, making their own API calls.
# Used for sites where robots.txt blocks scrapers but APIs are available.


class WikipediaFetcher:
    """Fetch Wikipedia articles via the REST API.

    Uses /api/rest_v1/page/summary/ for concise content, or
    /api/rest_v1/page/mobile-html/ for full articles.
    Bypasses robots.txt entirely since this is the official API.
    """

    # All Wikimedia sites
    _DOMAINS = {"wikipedia.org", "wiktionary.org", "wikimedia.org", "wikibooks.org"}

    def can_handle_url(self, url: str) -> bool:
        netloc = urlparse(url).netloc.lower()
        return any(netloc.endswith(d) for d in self._DOMAINS)

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        parsed = urlparse(url)
        path_parts = [p for p in parsed.path.split("/") if p]

        # Extract article title from /wiki/Article_Title
        if len(path_parts) < 2 or path_parts[0] != "wiki":
            return None
        article_title = path_parts[1]

        # Determine the language subdomain (en, de, fr, etc.)
        lang_host = parsed.netloc  # e.g. en.wikipedia.org

        # Use the summary API for a concise but useful result
        api_url = f"https://{lang_host}/api/rest_v1/page/summary/{article_title}"

        try:
            resp = requests.get(api_url, timeout=timeout, headers={
                "User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)",
                "Accept": "application/json",
            })
            resp.raise_for_status()
            data = resp.json()

            title = data.get("title", article_title.replace("_", " "))
            extract = data.get("extract", "")
            description = data.get("description", "")

            if not extract:
                return None

            # Also fetch full article HTML for deeper content
            full_content = self._fetch_full_article(lang_host, article_title, timeout)

            parts = []
            if description:
                parts.append(f"*{description}*\n")
            if full_content:
                parts.append(full_content)
            else:
                parts.append(extract)

            content = "\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="wikipedia",
            )

        except Exception as e:
            logger.error(f"Wikipedia API failed for {url}: {e}")
            return None

    def _fetch_full_article(self, host: str, title: str, timeout: int) -> str:
        """Fetch full article content via the mobile-html endpoint."""
        api_url = f"https://{host}/api/rest_v1/page/mobile-html/{title}"
        try:
            resp = requests.get(api_url, timeout=timeout, headers={
                "User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)",
                "Accept": "text/html",
            })
            resp.raise_for_status()

            if not BS4_AVAILABLE:
                return ""

            soup = BeautifulSoup(resp.content.decode("utf-8", errors="replace"), "html.parser")

            # Remove edit links, references section noise, navboxes
            for el in soup.select(".mw-ref, .navbox, .sistersitebox, .noprint, .mw-empty-elt"):
                el.decompose()

            # Extract section-by-section
            extractor = GenericExtractor()
            sections = soup.select("section")
            if sections:
                parts = []
                for section in sections:
                    text = extractor._html_to_markdown(section)
                    if text and len(text.strip()) > 20:
                        parts.append(text)
                return "\n\n".join(parts)

            return extractor._html_to_markdown(soup.body or soup)

        except Exception as e:
            logger.debug(f"Wikipedia full article fetch failed: {e}")
            return ""


class YouTubeFetcher:
    """Extract YouTube video metadata via oembed API.

    YouTube blocks scrapers and renders with JS, but the oembed
    endpoint returns title, author, thumbnail. We also parse the
    meta tags from the HTML response for description.
    """

    def can_handle_url(self, url: str) -> bool:
        netloc = urlparse(url).netloc.lower()
        return any(d in netloc for d in ("youtube.com", "youtu.be"))

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        # Normalize youtu.be short URLs
        parsed = urlparse(url)
        if "youtu.be" in parsed.netloc:
            video_id = parsed.path.lstrip("/")
            url = f"https://www.youtube.com/watch?v={video_id}"

        try:
            # oembed endpoint — always works, no auth needed
            oembed_url = f"https://www.youtube.com/oembed?url={url}&format=json"
            resp = requests.get(oembed_url, timeout=timeout)
            resp.raise_for_status()
            data = resp.json()

            title = data.get("title", "YouTube Video")
            author = data.get("author_name", "")
            author_url = data.get("author_url", "")

            # Try to get description from the page meta tags
            description = self._fetch_description(url, timeout)

            parts = []
            if author:
                parts.append(f"**Channel:** [{author}]({author_url})" if author_url else f"**Channel:** {author}")
            parts.append(f"**URL:** {url}")
            if description:
                parts.append(f"\n## Description\n\n{description}")

            content = "\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[author_url] if author_url else [],
                source_type="youtube",
            )

        except Exception as e:
            logger.error(f"YouTube oembed failed for {url}: {e}")
            return None

    def _fetch_description(self, url: str, timeout: int) -> str:
        """Fetch video description from page meta tags."""
        try:
            resp = requests.get(url, timeout=timeout, headers={
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/120.0.0.0 Safari/537.36"
                ),
            })
            resp.raise_for_status()

            if not BS4_AVAILABLE:
                # Regex fallback for description meta tag
                match = re.search(
                    r'<meta\s+name="description"\s+content="([^"]*)"', resp.text
                )
                return match.group(1) if match else ""

            soup = BeautifulSoup(resp.text, "html.parser")
            desc_tag = soup.find("meta", attrs={"name": "description"})
            if desc_tag and desc_tag.get("content"):
                return desc_tag["content"]

            # Try og:description
            og_desc = soup.find("meta", property="og:description")
            if og_desc and og_desc.get("content"):
                return og_desc["content"]

            return ""
        except Exception:
            return ""


class StackExchangeFetcher:
    """Fetch StackExchange Q&A via the public API v2.3.

    Covers stackoverflow.com, *.stackexchange.com, askubuntu.com,
    superuser.com, serverfault.com. Extracts question + top answers
    with vote counts.
    """

    _DOMAINS = {
        "stackoverflow.com", "superuser.com", "serverfault.com",
        "askubuntu.com", "mathoverflow.net",
    }

    def can_handle_url(self, url: str) -> bool:
        netloc = urlparse(url).netloc.lower()
        if netloc.startswith("www."):
            netloc = netloc[4:]
        return netloc in self._DOMAINS or netloc.endswith(".stackexchange.com")

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        parsed = urlparse(url)
        path_parts = [p for p in parsed.path.split("/") if p]

        # Extract question ID from /questions/12345/title-slug
        if len(path_parts) < 2 or path_parts[0] != "questions":
            return None
        try:
            question_id = int(path_parts[1])
        except ValueError:
            return None

        # Determine the API site parameter
        site = self._get_site_param(parsed.netloc)

        try:
            # Fetch question + answers in one call
            api_url = (
                f"https://api.stackexchange.com/2.3/questions/{question_id}"
                f"?order=desc&sort=votes&site={site}"
                f"&filter=withbody"  # Include body HTML
            )
            resp = requests.get(api_url, timeout=timeout)
            resp.raise_for_status()
            data = resp.json()

            items = data.get("items", [])
            if not items:
                return None

            q = items[0]
            title = q.get("title", "")
            q_body = self._html_to_text(q.get("body", ""))
            tags = q.get("tags", [])
            score = q.get("score", 0)
            answer_count = q.get("answer_count", 0)

            parts = [
                f"**Score:** {score}  |  **Answers:** {answer_count}",
            ]
            if tags:
                parts.append(f"**Tags:** {', '.join(tags)}")
            parts.append(f"\n## Question\n\n{q_body}")

            # Fetch answers
            answers = self._fetch_answers(question_id, site, timeout)
            for i, ans in enumerate(answers[:5], 1):
                accepted = " (Accepted)" if ans.get("is_accepted") else ""
                ans_body = self._html_to_text(ans.get("body", ""))
                parts.append(f"\n## Answer {i}{accepted} (Score: {ans.get('score', 0)})\n\n{ans_body}")

            content = "\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="stackexchange",
            )

        except Exception as e:
            logger.error(f"StackExchange API failed for {url}: {e}")
            return None

    def _fetch_answers(self, question_id: int, site: str, timeout: int) -> list:
        """Fetch answers sorted by votes."""
        try:
            api_url = (
                f"https://api.stackexchange.com/2.3/questions/{question_id}/answers"
                f"?order=desc&sort=votes&site={site}"
                f"&filter=withbody"
            )
            resp = requests.get(api_url, timeout=timeout)
            resp.raise_for_status()
            return resp.json().get("items", [])
        except Exception:
            return []

    def _get_site_param(self, netloc: str) -> str:
        """Convert domain to StackExchange API site parameter."""
        netloc = netloc.lower()
        if netloc.startswith("www."):
            netloc = netloc[4:]
        if netloc == "stackoverflow.com":
            return "stackoverflow"
        if netloc.endswith(".stackexchange.com"):
            return netloc.replace(".stackexchange.com", "")
        # askubuntu, superuser, serverfault, mathoverflow
        return netloc.replace(".com", "").replace(".net", "")

    def _html_to_text(self, html: str) -> str:
        """Convert SE answer HTML to clean markdown."""
        if not html:
            return ""
        if BS4_AVAILABLE:
            extractor = GenericExtractor()
            # Wrap in a div so _html_to_markdown can recurse into children
            soup = BeautifulSoup(f"<div>{html}</div>", "html.parser")
            return extractor._html_to_markdown(soup.div)
        # Minimal fallback: strip tags
        return re.sub(r"<[^>]+>", "", html).strip()


class RedditFetcher:
    """Fetch Reddit posts via the JSON API.

    Appending .json to any Reddit URL returns structured data.
    No authentication needed. Bypasses robots.txt block.
    """

    def can_handle_url(self, url: str) -> bool:
        netloc = urlparse(url).netloc.lower()
        return any(d in netloc for d in ("reddit.com", "redd.it"))

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        # Normalize URL and append .json
        parsed = urlparse(url)
        json_url = url.rstrip("/") + ".json"
        if ".json.json" in json_url:
            json_url = url  # Already has .json

        try:
            resp = requests.get(json_url, timeout=timeout, headers={
                "User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)",
            })
            resp.raise_for_status()
            data = resp.json()

            # Reddit returns a list: [post_listing, comments_listing]
            if not isinstance(data, list) or len(data) < 1:
                return None

            post_data = data[0].get("data", {}).get("children", [])
            if not post_data:
                return None

            post = post_data[0].get("data", {})
            title = post.get("title", "")
            selftext = post.get("selftext", "")
            subreddit = post.get("subreddit_name_prefixed", "")
            score = post.get("score", 0)
            author = post.get("author", "")
            num_comments = post.get("num_comments", 0)
            post_url = post.get("url", url)

            parts = [
                f"**{subreddit}**  |  Score: {score}  |  {num_comments} comments",
                f"**Posted by:** u/{author}",
            ]

            # If it's a link post, include the linked URL
            if post_url and post_url != url and not post_url.startswith("https://www.reddit.com"):
                parts.append(f"**Link:** {post_url}")

            if selftext:
                parts.append(f"\n{selftext}")

            # Extract top comments
            if len(data) > 1:
                comments = data[1].get("data", {}).get("children", [])
                top_comments = [c for c in comments[:5] if c.get("kind") == "t1"]
                if top_comments:
                    parts.append("\n## Top Comments\n")
                    for c in top_comments:
                        cd = c.get("data", {})
                        c_author = cd.get("author", "[deleted]")
                        c_body = cd.get("body", "")
                        c_score = cd.get("score", 0)
                        if c_body:
                            parts.append(f"**u/{c_author}** (Score: {c_score}):\n> {c_body}\n")

            content = "\n".join(parts)
            word_count = len(content.split())

            if not title:
                return None

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[post_url] if post_url != url else [],
                source_type="reddit",
            )

        except Exception as e:
            logger.error(f"Reddit JSON API failed for {url}: {e}")
            return None


class DevToFetcher:
    """Fetch Dev.to articles via the public API.

    Dev.to blocks scrapers via robots.txt but has a public REST API
    that returns articles as markdown.
    """

    def can_handle_url(self, url: str) -> bool:
        netloc = urlparse(url).netloc.lower()
        if netloc.startswith("www."):
            netloc = netloc[4:]
        return netloc == "dev.to"

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        parsed = urlparse(url)
        path_parts = [p for p in parsed.path.split("/") if p]

        # Need at least /username/article-slug
        if len(path_parts) < 2:
            return None
        # Skip tag pages (/t/python), top pages, etc.
        if path_parts[0] in ("t", "top", "latest", "search", "tags", "settings"):
            return None

        try:
            username = path_parts[0]
            slug = path_parts[1]
            headers = {"User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)"}

            # Fetch article directly by username/slug path
            resp = requests.get(
                f"https://dev.to/api/articles/{username}/{slug}",
                timeout=timeout, headers=headers,
            )
            resp.raise_for_status()
            article = resp.json()

            if not article or not article.get("title"):
                return None

            title = article.get("title", "")
            body = article.get("body_markdown", "")

            # Fallback to body_html → markdown
            if not body:
                body_html = article.get("body_html", "")
                if body_html and BS4_AVAILABLE:
                    soup = BeautifulSoup(body_html, "html.parser")
                    body = GenericExtractor()._html_to_markdown(soup)

            if not body or not title:
                return None

            # Build metadata header
            author = article.get("user", {}).get("name", "")
            tags = article.get("tags", article.get("tag_list", []))
            if isinstance(tags, str):
                tags = [t.strip() for t in tags.split(",") if t.strip()]
            reactions = article.get("positive_reactions_count", 0)
            date = article.get("readable_publish_date", "")

            parts = []
            if author:
                parts.append(f"**Author:** {author}")
            if date:
                parts.append(f"**Published:** {date}")
            if tags:
                parts.append(f"**Tags:** {', '.join(tags)}")
            if reactions:
                parts.append(f"**Reactions:** {reactions}")
            parts.append("")
            parts.append(body)

            content = "\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="devto",
            )

        except Exception as e:
            logger.error(f"Dev.to API failed for {url}: {e}")
            return None


class FandomFetcher:
    """Fetch Fandom wiki articles via the MediaWiki API.

    All *.fandom.com wikis use MediaWiki and block scrapers via robots.txt,
    but the API is accessible. Returns wikitext converted to markdown.
    """

    def can_handle_url(self, url: str) -> bool:
        return urlparse(url).netloc.lower().endswith(".fandom.com")

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        parsed = urlparse(url)
        netloc = parsed.netloc.lower()
        wiki = netloc.split(".fandom.com")[0]

        path_parts = [p for p in parsed.path.split("/") if p]
        if len(path_parts) < 2 or path_parts[0] != "wiki":
            return None

        from urllib.parse import unquote
        page_title = unquote("/".join(path_parts[1:]))

        try:
            resp = requests.get(
                f"https://{wiki}.fandom.com/api.php",
                params={
                    "action": "parse",
                    "page": page_title,
                    "format": "json",
                    "prop": "wikitext|categories|displaytitle",
                    "redirects": "1",
                },
                timeout=timeout,
                headers={"User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)"},
            )
            resp.raise_for_status()
            data = resp.json()

            if "error" in data:
                return None

            parse = data.get("parse", {})
            title = parse.get("displaytitle", parse.get("title", page_title.replace("_", " ")))
            # Strip HTML from displaytitle
            if "<" in title and BS4_AVAILABLE:
                title = BeautifulSoup(title, "html.parser").get_text()

            wikitext = parse.get("wikitext", {}).get("*", "")
            if not wikitext:
                return None

            content = self._wikitext_to_markdown(wikitext)

            # Add categories
            categories = parse.get("categories", [])
            if categories:
                cat_names = [
                    c.get("*", "").replace("_", " ") for c in categories
                    if not c.get("hidden")
                ]
                if cat_names:
                    content += f"\n\n**Categories:** {', '.join(cat_names)}"

            word_count = len(content.split())
            if word_count < 10:
                return None

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="fandom",
            )

        except Exception as e:
            logger.error(f"Fandom API failed for {url}: {e}")
            return None

    @staticmethod
    def _wikitext_to_markdown(wikitext: str) -> str:
        """Convert MediaWiki wikitext to clean markdown."""
        text = wikitext

        # Strip templates {{...}} (infoboxes, navboxes — mostly noise)
        # Handle nested templates by iterating
        for _ in range(5):
            prev = text
            text = re.sub(r"\{\{[^{}]*\}\}", "", text)
            if text == prev:
                break

        # Strip references <ref>...</ref> and <ref ... />
        text = re.sub(r"<ref[^>]*>.*?</ref>", "", text, flags=re.DOTALL)
        text = re.sub(r"<ref[^/]*/\s*>", "", text)

        # Strip HTML comments
        text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)

        # Strip remaining HTML tags
        text = re.sub(r"<[^>]+>", "", text)

        # Headings: == H2 == → ## H2
        text = re.sub(r"^=====\s*(.+?)\s*=====", r"##### \1", text, flags=re.MULTILINE)
        text = re.sub(r"^====\s*(.+?)\s*====", r"#### \1", text, flags=re.MULTILINE)
        text = re.sub(r"^===\s*(.+?)\s*===", r"### \1", text, flags=re.MULTILINE)
        text = re.sub(r"^==\s*(.+?)\s*==", r"## \1", text, flags=re.MULTILINE)

        # Bold/italic: '''bold''' → **bold**, ''italic'' → *italic*
        # Bold MUST be processed before italic (''' before '')
        text = re.sub(r"'{3}(.+?)'{3}", r"**\1**", text)
        text = re.sub(r"'{2}(.+?)'{2}", r"*\1*", text)

        # Links: [[Page|Text]] → Text, [[Page]] → Page
        text = re.sub(r"\[\[[^|\]]+\|([^\]]+)\]\]", r"\1", text)
        text = re.sub(r"\[\[([^\]]+)\]\]", r"\1", text)

        # External links: [http://url text] → text
        text = re.sub(r"\[https?://\S+\s+([^\]]+)\]", r"\1", text)
        text = re.sub(r"\[https?://\S+\]", "", text)

        # Lists: * item → - item (but not ** bold markers)
        text = re.sub(r"^\*\s+", "- ", text, flags=re.MULTILINE)

        # Clean excessive whitespace
        text = re.sub(r"\n{3,}", "\n\n", text)

        return text.strip()


class SemanticScholarFetcher:
    """Fetch academic paper metadata via the Semantic Scholar API.

    Resolves papers from multiple academic domains: DOI, PubMed,
    Semantic Scholar, ACM, IEEE, ScienceDirect, ResearchGate,
    Springer, Nature. Returns title, abstract, authors, citations,
    AI-generated TLDR, and open access PDF links.
    """

    _ACADEMIC_DOMAINS = {
        "semanticscholar.org",
        "doi.org", "dx.doi.org",
        "dl.acm.org", "acm.org",
        "ieeexplore.ieee.org", "ieee.org",
        "sciencedirect.com",
        "researchgate.net",
        "springer.com", "link.springer.com",
        "nature.com",
        "pubmed.ncbi.nlm.nih.gov",
        "ncbi.nlm.nih.gov",
        "wiley.com", "onlinelibrary.wiley.com",
        "plos.org", "journals.plos.org",
    }

    _S2_API = "https://api.semanticscholar.org/graph/v1/paper"
    _S2_FIELDS = "title,abstract,authors,year,venue,citationCount,referenceCount,tldr,externalIds,openAccessPdf"

    def can_handle_url(self, url: str) -> bool:
        netloc = urlparse(url).netloc.lower()
        if netloc.startswith("www."):
            netloc = netloc[4:]
        return any(netloc == d or netloc.endswith("." + d) for d in self._ACADEMIC_DOMAINS)

    def fetch_and_extract(self, url: str, timeout: int = 15) -> Optional[ScrapedPage]:
        paper_id = self._resolve_paper_id(url)
        if not paper_id:
            return None

        try:
            from urllib.parse import quote
            api_url = f"{self._S2_API}/{quote(paper_id, safe=':')}?fields={self._S2_FIELDS}"

            resp = requests.get(
                api_url, timeout=timeout,
                headers={"User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)"},
            )

            if resp.status_code == 404:
                return None
            if resp.status_code == 429:
                # Rate limited — respect Retry-After or wait 5s, retry once
                import time
                wait = int(resp.headers.get("Retry-After", 5))
                time.sleep(min(wait, 10))
                resp = requests.get(
                    api_url, timeout=timeout,
                    headers={"User-Agent": "FSS-Mini-RAG-Research/2.2 (research tool)"},
                )
            resp.raise_for_status()
            data = resp.json()

            title = data.get("title", "")
            if not title:
                return None

            abstract = data.get("abstract", "")

            # Build rich content
            parts = []

            # Authors
            authors = data.get("authors", [])
            if authors:
                author_names = [a.get("name", "") for a in authors if a.get("name")]
                if author_names:
                    parts.append(f"**Authors:** {', '.join(author_names)}")

            # Year & venue
            year = data.get("year")
            venue = data.get("venue", "")
            if year and venue:
                parts.append(f"**Published:** {venue}, {year}")
            elif year:
                parts.append(f"**Year:** {year}")

            # Citation stats
            citations = data.get("citationCount", 0)
            references = data.get("referenceCount", 0)
            if citations or references:
                parts.append(f"**Citations:** {citations}  |  **References:** {references}")

            # External IDs
            ext_ids = data.get("externalIds", {})
            id_parts = []
            if ext_ids.get("DOI"):
                id_parts.append(f"DOI: {ext_ids['DOI']}")
            if ext_ids.get("ArXiv"):
                id_parts.append(f"arXiv: {ext_ids['ArXiv']}")
            if ext_ids.get("PMID"):
                id_parts.append(f"PMID: {ext_ids['PMID']}")
            if id_parts:
                parts.append(f"**IDs:** {' | '.join(id_parts)}")

            # Open access PDF
            oa_pdf = data.get("openAccessPdf")
            if oa_pdf and isinstance(oa_pdf, dict) and oa_pdf.get("url"):
                parts.append(f"**PDF:** [{oa_pdf['url']}]({oa_pdf['url']})")

            parts.append("")

            # TLDR (AI-generated summary)
            tldr = data.get("tldr")
            if tldr and isinstance(tldr, dict) and tldr.get("text"):
                parts.append(f"## TL;DR\n\n{tldr['text']}")
                parts.append("")

            # Abstract
            if abstract:
                parts.append(f"## Abstract\n\n{abstract}")

            content = "\n".join(parts)
            word_count = len(content.split())

            return ScrapedPage(
                url=url,
                title=title,
                content=content,
                scraped_at=datetime.now().isoformat(),
                word_count=word_count,
                links=[],
                source_type="academic",
            )

        except Exception as e:
            logger.error(f"Semantic Scholar API failed for {url}: {e}")
            return None

    def _resolve_paper_id(self, url: str) -> Optional[str]:
        """Convert a URL to a Semantic Scholar paper identifier."""
        parsed = urlparse(url)
        netloc = parsed.netloc.lower()
        if netloc.startswith("www."):
            netloc = netloc[4:]
        path = parsed.path

        # Direct S2 URL: semanticscholar.org/paper/Title/hexid
        if "semanticscholar.org" in netloc:
            parts = [p for p in path.split("/") if p]
            if len(parts) >= 2 and parts[0] == "paper":
                return parts[-1]
            return None

        # DOI URLs: doi.org/10.xxxx/yyyy
        if netloc in ("doi.org", "dx.doi.org"):
            doi = path.lstrip("/")
            if not doi:
                return None
            # arXiv DOIs (10.48550/arXiv.*) resolve better via ARXIV: prefix
            if doi.startswith("10.48550/arXiv."):
                arxiv_id = doi.replace("10.48550/arXiv.", "")
                return f"ARXIV:{arxiv_id}"
            return f"DOI:{doi}"

        # PubMed: pubmed.ncbi.nlm.nih.gov/12345678
        if "pubmed" in netloc or "ncbi.nlm.nih.gov" in netloc:
            parts = [p for p in path.split("/") if p]
            for part in reversed(parts):
                if part.isdigit():
                    return f"PMID:{part}"
            return None

        # For other academic domains, try URL-based resolution
        return f"URL:{url}"


# Registry of all extractors, ordered by specificity (most specific first)
_EXTRACTORS: List[ContentExtractor] = [
    ArxivExtractor(),
    GitHubExtractor(),
    PDFExtractor(),
    DocxExtractor(),
    SpreadsheetExtractor(),
    PptxExtractor(),
    EpubExtractor(),
    RtfExtractor(),
    RssFeedExtractor(),
    MarkdownPassthroughExtractor(),
    GenericExtractor(),  # Fallback — must be last
]

# API-based direct fetchers — checked before the normal fetch pipeline
_DIRECT_FETCHERS: List[DirectFetcher] = [
    WikipediaFetcher(),
    YouTubeFetcher(),
    StackExchangeFetcher(),
    RedditFetcher(),
    DevToFetcher(),
    FandomFetcher(),
    SemanticScholarFetcher(),
]


def get_direct_fetcher(url: str) -> Optional[DirectFetcher]:
    """Find a direct fetcher for a URL, if one exists."""
    for fetcher in _DIRECT_FETCHERS:
        if fetcher.can_handle_url(url):
            return fetcher
    return None


def get_extractor(url: str, content_type: str) -> Optional[ContentExtractor]:
    """Find the best extractor for a given URL and content type."""
    for extractor in _EXTRACTORS:
        if extractor.can_handle(url, content_type):
            return extractor
    return None


def extract_content(url: str, raw: bytes, content_type: str) -> Optional[ScrapedPage]:
    """Extract content from raw bytes using the appropriate extractor."""
    from .scrape_registry import log_scrape

    extractor = get_extractor(url, content_type)
    if extractor is None:
        logger.warning(f"No extractor available for {content_type} from {url}")
        log_scrape(
            url=url, extractor_name="none", success=False,
            error=f"No extractor for {content_type}", content_length=len(raw),
        )
        return None

    extractor_name = type(extractor).__name__
    page = extractor.extract(url, raw, content_type)

    if page:
        log_scrape(
            url=url, extractor_name=extractor_name, success=True,
            word_count=page.word_count, title=page.title,
            source_type=page.source_type, content_length=len(raw),
            has_main_content=page.word_count > 100,
            doc_metadata=page.metadata if page.metadata else None,
        )
    else:
        log_scrape(
            url=url, extractor_name=extractor_name, success=False,
            error="Extraction returned None", content_length=len(raw),
        )

    return page
